from __future__ import annotations
import os
import asyncio
from pathlib import Path
from typing import List, Dict, Any

import fitz  # PyMuPDF
from dataflow_agent.workflow.registry import register
from dataflow_agent.graphbuilder.graph_builder import GenericGraphBuilder
from dataflow_agent.logger import get_logger
from dataflow_agent.state import KBPodcastState, MainState
from dataflow_agent.agentroles import create_agent
from dataflow_agent.utils import get_project_root
import re
import wave
from dataflow_agent.toolkits.multimodaltool.req_tts import (
    generate_speech_bytes_async,
    split_tts_text
)

log = get_logger(__name__)

# Try importing office libraries
try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

@register("kb_podcast")
def create_kb_podcast_graph() -> GenericGraphBuilder:
    """
    Workflow for Knowledge Base Podcast Generation
    Steps:
    1. Parse uploaded files (PDF/Office)
    2. Generate podcast script using LLM
    3. Generate audio using TTS
    """
    builder = GenericGraphBuilder(state_model=KBPodcastState, entry_point="_start_")

    def _extract_text_result(state: MainState, role_name: str) -> str:
        try:
            result = state.agent_results.get(role_name, {}).get("results", {})
            if isinstance(result, dict):
                return result.get("text") or result.get("raw") or ""
            if isinstance(result, str):
                return result
        except Exception:
            return ""
        return ""

    def _start_(state: KBPodcastState) -> KBPodcastState:
        # Ensure request fields
        if not state.request.files:
            state.request.files = []

        # Initialize output directory
        if not state.result_path:
            project_root = get_project_root()
            import time
            ts = int(time.time())
            email = getattr(state.request, 'email', 'default')
            output_dir = project_root / "outputs" / "kb_outputs" / email / f"{ts}_podcast"
            output_dir.mkdir(parents=True, exist_ok=True)
            state.result_path = str(output_dir)
        else:
            Path(state.result_path).mkdir(parents=True, exist_ok=True)

        state.file_contents = []
        state.podcast_script = ""
        state.audio_path = ""
        return state

    async def parse_files_node(state: KBPodcastState) -> KBPodcastState:
        """
        Parse all files and extract content
        """
        files = state.request.files
        if not files:
            state.file_contents = []
            return state

        async def process_file(file_path: str) -> Dict[str, Any]:
            file_path_obj = Path(file_path)
            filename = file_path_obj.name

            if not file_path_obj.exists():
                return {
                    "filename": filename,
                    "content": f"[Error: File not found {file_path}]"
                }

            suffix = file_path_obj.suffix.lower()
            raw_content = ""

            try:
                # PDF
                if suffix == ".pdf":
                    try:
                        doc = fitz.open(file_path)
                        text = ""
                        for page in doc:
                            text += page.get_text() + "\n"
                        raw_content = text
                    except Exception as e:
                        raw_content = f"[Error parsing PDF: {e}]"

                # Word
                elif suffix in [".docx", ".doc"]:
                    if Document is None:
                         raw_content = "[Error: python-docx not installed]"
                    else:
                        try:
                            doc = Document(file_path)
                            raw_content = "\n".join([p.text for p in doc.paragraphs])
                        except Exception as e:
                             raw_content = f"[Error parsing Docx: {e}]"

                # PPT
                elif suffix in [".pptx", ".ppt"]:
                    if Presentation is None:
                        raw_content = "[Error: python-pptx not installed]"
                    else:
                        try:
                            prs = Presentation(file_path)
                            text = ""
                            for i, slide in enumerate(prs.slides):
                                text += f"--- Slide {i+1} ---\n"
                                for shape in slide.shapes:
                                    if hasattr(shape, "text"):
                                        text += shape.text + "\n"
                            raw_content = text
                        except Exception as e:
                            raw_content = f"[Error parsing PPT: {e}]"

                else:
                    try:
                        with open(file_path, "r", encoding="utf-8") as f:
                            raw_content = f.read()
                    except:
                        raw_content = "[Unsupported file type]"

            except Exception as e:
                 raw_content = f"[Parse Error: {e}]"

            # Truncate content
            truncated_content = raw_content[:50000] if len(raw_content) > 50000 else raw_content

            return {
                "filename": filename,
                "content": truncated_content
            }

        # Run in parallel
        tasks = [process_file(f) for f in files]
        results = await asyncio.gather(*tasks)

        state.file_contents = results
        return state

    async def generate_script_node(state: KBPodcastState) -> KBPodcastState:
        """
        Generate podcast script using LLM
        """
        if not state.file_contents:
            state.podcast_script = "No content available for podcast generation."
            return state

        # Format file contents
        contents_str = ""
        for item in state.file_contents:
            contents_str += f"=== {item['filename']} ===\n{item['content']}\n\n"

        # Podcast script prompt
        language = state.request.language
        mode = getattr(state.request, "podcast_mode", "monologue")
        if mode == "dialog":
            speaker_a = "主持人" if language == "zh" else "Host"
            speaker_b = "嘉宾" if language == "zh" else "Guest"
            prompt = f"""你是一位专业的知识播客制作人。基于以下资料，生成一段5-10分钟的双人对话播客脚本。

要求：
1. 口语化、生动有趣，避免书面语
2. 结构清晰：开场白 → 核心内容 → 总结
3. 使用类比和例子帮助理解
4. 适当加入互动性语言（"你可能会想..."）
5. 使用{language}语言
6. 严格使用如下格式逐行输出（每行一个角色）：
{speaker_a}: ...
{speaker_b}: ...

资料内容：
{contents_str}

请生成播客脚本："""
        else:
            prompt = f"""你是一位专业的知识播客主播。基于以下资料，生成一段5-10分钟的知识播客脚本。

要求：
1. 口语化、生动有趣，避免书面语
2. 结构清晰：开场白 → 核心内容 → 总结
3. 使用类比和例子帮助理解
4. 适当加入互动性语言（"你可能会想..."）
5. 使用{language}语言

资料内容：
{contents_str}

请生成播客脚本："""

        try:
            agent = create_agent(
                name="kb_prompt_agent",
                model_name=state.request.model,
                chat_api_url=state.request.chat_api_url,
                temperature=0.7,
                parser_type="text"
            )

            temp_state = MainState(request=state.request)
            res_state = await agent.execute(temp_state, prompt=prompt)

            state.podcast_script = _extract_text_result(res_state, "kb_prompt_agent") or "[Script generation failed]"
        except Exception as e:
            log.error(f"Script generation failed: {e}")
            state.podcast_script = f"[Script generation error: {e}]"

        # Save script to file
        try:
            script_path = Path(state.result_path) / "script.txt"
            script_path.write_text(state.podcast_script, encoding="utf-8")
        except Exception as e:
            log.error(f"Failed to save script: {e}")

        return state

    async def generate_audio_node(state: KBPodcastState) -> KBPodcastState:
        """
        Generate audio using TTS
        """
        if not state.podcast_script or state.podcast_script.startswith("["):
            state.audio_path = ""
            return state

        try:
            audio_path = str(Path(state.result_path) / "podcast.wav")
            mode = getattr(state.request, "podcast_mode", "monologue")
            max_chars = 1500
            concurrency = 4

            segments = []
            if mode == "dialog":
                language = state.request.language
                speaker_a = "主持人" if language == "zh" else "Host"
                speaker_b = "嘉宾" if language == "zh" else "Guest"
                speaker_map = {
                    speaker_a.lower(): "A",
                    speaker_b.lower(): "B",
                    "a": "A",
                    "b": "B",
                    "speaker a": "A",
                    "speaker b": "B",
                    "角色a": "A",
                    "角色b": "B",
                    "主播": "A",
                    "嘉宾": "B",
                }
                pattern = re.compile(r"^\s*([^:：]{1,20})\s*[:：]\s*(.+)$")
                current_speaker = "A"
                for raw_line in state.podcast_script.splitlines():
                    line = raw_line.strip()
                    if not line:
                        continue
                    m = pattern.match(line)
                    if m:
                        label = m.group(1).strip().lower()
                        content = m.group(2).strip()
                        mapped = speaker_map.get(label)
                        if mapped:
                            current_speaker = mapped
                        if content:
                            segments.append({"speaker": current_speaker, "text": content})
                        continue
                    # No label, append to current speaker
                    if segments and segments[-1]["speaker"] == current_speaker:
                        segments[-1]["text"] = f"{segments[-1]['text']} {line}"
                    else:
                        segments.append({"speaker": current_speaker, "text": line})

                expanded = []
                for seg in segments:
                    for chunk in split_tts_text(seg["text"], max_chars):
                        expanded.append({
                            "speaker": seg["speaker"],
                            "text": chunk
                        })
                segments = expanded
            else:
                for chunk in split_tts_text(state.podcast_script, max_chars):
                    segments.append({"speaker": "A", "text": chunk})

            if not segments:
                raise RuntimeError("No valid TTS segments generated from script")

            sem = asyncio.Semaphore(concurrency)

            async def _run(seg):
                voice = state.request.voice_name if seg["speaker"] == "A" else state.request.voice_name_b
                async with sem:
                    return await generate_speech_bytes_async(
                        text=seg["text"],
                        api_url=state.request.chat_api_url,
                        api_key=state.request.api_key,
                        model=state.request.tts_model,
                        voice_name=voice,
                    )

            async def _run_no_sem(seg):
                voice = state.request.voice_name if seg["speaker"] == "A" else state.request.voice_name_b
                return await generate_speech_bytes_async(
                    text=seg["text"],
                    api_url=state.request.chat_api_url,
                    api_key=state.request.api_key,
                    model=state.request.tts_model,
                    voice_name=voice,
                )

            async def _run_with_retry(seg, attempts=3, base_delay=0.8, use_sem=True):
                last_err = None
                for i in range(attempts):
                    try:
                        if use_sem:
                            return await _run(seg)
                        return await _run_no_sem(seg)
                    except Exception as e:
                        last_err = e
                        if i < attempts - 1:
                            await asyncio.sleep(base_delay * (i + 1))
                        continue
                raise last_err

            tasks = [asyncio.create_task(_run_with_retry(seg, attempts=2, use_sem=True)) for seg in segments]
            results = await asyncio.gather(*tasks, return_exceptions=True)

            failed_indices = [i for i, r in enumerate(results) if isinstance(r, Exception)]
            if failed_indices:
                log.warning(f"TTS retry sequentially for {len(failed_indices)} failed segment(s)")
                for i in failed_indices:
                    results[i] = await _run_with_retry(segments[i], attempts=3, use_sem=False)

            audio_chunks = results

            os.makedirs(os.path.dirname(os.path.abspath(audio_path)), exist_ok=True)
            with wave.open(audio_path, "wb") as wav_file:
                wav_file.setnchannels(1)        # 1 Channel
                wav_file.setsampwidth(2)        # 16 bit = 2 bytes
                wav_file.setframerate(24000)    # 24kHz
                wav_file.writeframes(b"".join(audio_chunks))

            state.audio_path = audio_path
            log.info(f"Audio generated successfully: {audio_path}")
        except Exception as e:
            log.error(f"Audio generation failed: {e}")
            err_str = str(e)
            if "503" in err_str or "model_not_found" in err_str or "model not found" in err_str.lower():
                state.audio_path = (
                    "[TTS 模型不可用：当前 API 不支持所选 TTS 模型（如 gemini-2.5-pro-preview-tts）。"
                    "请到「播客」设置中更换 TTS 模型，或使用支持该模型的 API 服务商。]"
                )
            else:
                state.audio_path = f"[Audio generation error: {e}]"

        return state

    nodes = {
        "_start_": _start_,
        "parse_files": parse_files_node,
        "generate_script": generate_script_node,
        "generate_audio": generate_audio_node,
        "_end_": lambda s: s
    }

    edges = [
        ("_start_", "parse_files"),
        ("parse_files", "generate_script"),
        ("generate_script", "generate_audio"),
        ("generate_audio", "_end_")
    ]

    builder.add_nodes(nodes).add_edges(edges)
    return builder
