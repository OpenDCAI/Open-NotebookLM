import asyncio
import json
import os
import re
import shutil
import subprocess
import time
from pathlib import Path
from urllib.parse import urlparse, unquote
from fastapi import APIRouter, UploadFile, File, Form, HTTPException, Body
from typing import Optional, List, Dict, Any

import fitz  # PyMuPDF

from dataflow_agent.state import IntelligentQARequest, IntelligentQAState, KBPodcastRequest, KBPodcastState, KBMindMapRequest, KBMindMapState
from dataflow_agent.workflow.wf_intelligent_qa import create_intelligent_qa_graph
from dataflow_agent.workflow.wf_kb_podcast import create_kb_podcast_graph
from dataflow_agent.workflow.wf_kb_mindmap import create_kb_mindmap_graph
from dataflow_agent.toolkits.ragtool.vector_store_tool import process_knowledge_base_files, VectorStoreManager
from dataflow_agent.utils import get_project_root
from dataflow_agent.logger import get_logger
from dataflow_agent.workflow import run_workflow

log = get_logger(__name__)
from fastapi_app.config import settings
from fastapi_app.schemas import Paper2PPTRequest
from fastapi_app.utils import _from_outputs_url, _to_outputs_url
from fastapi_app.workflow_adapters.wa_paper2ppt import _init_state_from_request
from fastapi_app.dependencies.auth import get_supabase_admin_client
from fastapi_app.notebook_paths import NotebookPaths, get_notebook_paths
from fastapi_app.source_manager import SourceManager
from fastapi_app.services.fast_research_service import fast_research_search
from fastapi_app.services.deep_research_report_service import generate_report_from_search
from dataflow_agent.toolkits.research_tools import fetch_page_text

router = APIRouter(prefix="/kb", tags=["Knowledge Base"])

# Link sources JSON filename under notebook dir (excluded from regular file list)
LINK_SOURCES_FILENAME = "link_sources.json"

# Base directory for storing KB files
# Layout: outputs/kb_data/{email}/{notebook_id}/ for per-notebook isolation
KB_BASE_DIR = Path("outputs/kb_data")
OUTPUTS_BASE = Path("outputs/kb_outputs")


def _notebook_dir(email: str, notebook_id: Optional[str]) -> Path:
    """User + notebook scoped dir under kb_data. Use raw email on disk so StaticFiles can resolve URL-decoded path."""
    root = get_project_root()
    base = root / KB_BASE_DIR / (email or "default")
    if notebook_id:
        return base / notebook_id.replace("/", "_").replace("\\", "_")[:128]
    return base / "_shared"


def _outputs_dir(email: str, notebook_id: Optional[str], subdir: str) -> Path:
    """User + notebook scoped output dir. Use raw email on disk for StaticFiles resolution."""
    root = get_project_root()
    base = root / OUTPUTS_BASE / (email or "default")
    if notebook_id:
        base = base / notebook_id.replace("/", "_").replace("\\", "_")[:128]
    else:
        base = base / "_shared"
    return base / subdir


def _get_cjk_font_path() -> Optional[str]:
    """返回系统中文字体路径，用于 PDF 内中文显示；无则返回 None。"""
    candidates = [
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/opentype/noto/NotoSerifCJK-Regular.ttc",
        "/usr/share/fonts/truetype/arphic/uming.ttc",
    ]
    for path in candidates:
        if Path(path).exists():
            return path
    return None


def _text_to_pdf(text: str, output_path: str) -> None:
    """将长文本生成为多页 PDF（PyMuPDF），不依赖 kb_page_content / paper2ppt workflow。支持中文（CJK 字体）。"""
    text = (text or "").strip()
    if not text:
        raise ValueError("Report text is empty")
    doc = fitz.open()
    rect = fitz.Rect(50, 50, 545, 802)
    fontsize = 11
    max_chars_per_page = 3200
    fontfile = _get_cjk_font_path()
    fontname = "notocjk" if fontfile else "helv"
    if fontfile:
        # 使用中文字体，否则中文会不显示
        pass
    paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
    if not paragraphs:
        paragraphs = [text]
    current: List[str] = []
    current_len = 0
    for p in paragraphs:
        need = len(p) + (2 if current else 0)
        if current_len + need > max_chars_per_page and current:
            page = doc.new_page(width=595, height=842)
            if fontfile:
                page.insert_textbox(rect, "\n\n".join(current), fontsize=fontsize, fontname=fontname, fontfile=fontfile)
            else:
                page.insert_textbox(rect, "\n\n".join(current), fontsize=fontsize, fontname=fontname)
            current = [p]
            current_len = len(p)
        else:
            current.append(p)
            current_len += need
    if current:
        page = doc.new_page(width=595, height=842)
        if fontfile:
            page.insert_textbox(rect, "\n\n".join(current), fontsize=fontsize, fontname=fontname, fontfile=fontfile)
        else:
            page.insert_textbox(rect, "\n\n".join(current), fontsize=fontsize, fontname=fontname)
    doc.save(output_path)
    doc.close()


ALLOWED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".png", ".jpg", ".jpeg", ".mp4", ".md"}

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg"}
DOC_EXTENSIONS = {".pdf", ".docx", ".doc", ".pptx", ".ppt", ".md", ".markdown"}


def _find_mineru_stem_dir(
    pdf_stem: str,
    email: str,
    notebook_id: Optional[str],
    notebook_title: Optional[str] = None,
) -> Optional[Path]:
    """
    查找指定 pdf_stem 的 MinerU 输出目录。
    查找顺序：
    1. 笔记本新布局: outputs/{title}_{id}/sources/{pdf_stem}/mineru/
    2. kb_mineru 新结构: kb_mineru/{email}/{notebook_id}/{pdf_stem}/auto/
    3. kb_mineru 旧结构: kb_mineru/{email}/{notebook_id}/{uuid}/{pdf_stem}/auto/
    返回包含 auto/ 或 hybrid_auto/ 的目录，找不到返回 None。
    """
    project_root = get_project_root()

    # 1) 笔记本新布局: outputs/{title}_{id}/sources/{pdf_stem}/mineru/
    if notebook_id:
        nb_paths = get_notebook_paths(notebook_id, notebook_title or "", email)
        mineru_dir = nb_paths.sources_dir / pdf_stem / "mineru"
        if mineru_dir.exists():
            # 直接在 mineru/ 下找 auto/ 或 hybrid_auto/
            for sub in ("auto", "hybrid_auto"):
                if (mineru_dir / sub).is_dir() and list((mineru_dir / sub).glob("*.md")):
                    log.info("[find_mineru] 在新布局找到缓存: %s/%s", mineru_dir, sub)
                    return mineru_dir
            # 兼容: mineru/{pdf_stem}/auto/ (MinerU 可能多嵌套一层)
            nested = mineru_dir / pdf_stem
            if nested.exists():
                for sub in ("auto", "hybrid_auto"):
                    if (nested / sub).is_dir() and list((nested / sub).glob("*.md")):
                        log.info("[find_mineru] 在新布局(嵌套)找到缓存: %s/%s", nested, sub)
                        return nested

    # 2) Legacy: kb_mineru/{email}/{notebook_id}/
    safe_nb = (notebook_id or "_shared").replace("/", "_").replace("\\", "_")[:128]
    mineru_base = project_root / "outputs" / "kb_mineru" / (email or "default") / safe_nb

    if not mineru_base.exists():
        return None

    stem_dir = mineru_base / pdf_stem
    if stem_dir.exists():
        for sub in ("auto", "hybrid_auto"):
            if (stem_dir / sub).is_dir() and list((stem_dir / sub).glob("*.md")):
                return stem_dir

    # 3) 旧结构兼容：kb_mineru/{email}/{nb}/{uuid}/{pdf_stem}/auto/
    for child in mineru_base.iterdir():
        if not child.is_dir():
            continue
        nested = child / pdf_stem
        if not nested.exists():
            continue
        for sub in ("auto", "hybrid_auto"):
            if (nested / sub).is_dir() and list((nested / sub).glob("*.md")):
                return nested

    return None


def _read_mineru_md_if_cached(
    pdf_path: Path,
    email: str,
    notebook_id: Optional[str],
    max_chars: int = 50000,
    notebook_title: Optional[str] = None,
) -> Optional[str]:
    """
    尝试从已有的 MinerU 缓存中读取 markdown 内容。
    找到则返回 markdown 文本，否则返回 None。
    """
    stem_dir = _find_mineru_stem_dir(pdf_path.stem, email, notebook_id, notebook_title)
    if stem_dir is None:
        return None

    for sub in ("auto", "hybrid_auto"):
        candidate = stem_dir / sub
        if not candidate.is_dir():
            continue
        md_files = list(candidate.glob("*.md"))
        if md_files:
            try:
                text = md_files[0].read_text(encoding="utf-8")
                if text.strip():
                    log.info("[read_mineru_md] 从缓存读取 %s, len=%s", md_files[0], len(text))
                    return text[:max_chars] if len(text) > max_chars else text
            except Exception as e:
                log.warning("[read_mineru_md] 读取失败 %s: %s", md_files[0], e)
    return None


def _reuse_mineru_cache(
    pdf_paths: List[Path],
    output_dir: Path,
    email: str,
    notebook_id: Optional[str],
    notebook_title: Optional[str] = None,
) -> int:
    """
    将已有的 MinerU 解析结果复制/软链到 PPT workflow 的 output_dir 下，
    使 parse_pdf_pages 能直接发现 {output_dir}/{pdf_stem}/auto/*.md 而跳过重新解析。
    返回成功复用的 PDF 数量。
    """
    reused = 0
    for pdf_path in pdf_paths:
        stem = pdf_path.stem
        cached_stem_dir = _find_mineru_stem_dir(stem, email, notebook_id, notebook_title)
        if cached_stem_dir is None:
            log.info("[reuse_mineru] 未找到 %s 的 MinerU 缓存", stem)
            continue

        target = output_dir / stem
        if target.exists():
            # 已存在（可能之前已复用或本次 workflow 已生成），跳过
            log.info("[reuse_mineru] 目标已存在，跳过: %s", target)
            reused += 1
            continue

        try:
            target.symlink_to(cached_stem_dir.resolve())
            log.info("[reuse_mineru] 软链成功: %s -> %s", target, cached_stem_dir)
            reused += 1
        except OSError:
            # 软链失败（跨文件系统等），回退到复制
            try:
                shutil.copytree(str(cached_stem_dir), str(target))
                log.info("[reuse_mineru] 复制成功: %s -> %s", target, cached_stem_dir)
                reused += 1
            except Exception as e:
                log.warning("[reuse_mineru] 复制失败 %s: %s", stem, e)

    return reused


def _resolve_local_path(path_or_url: str) -> Path:
    if not path_or_url:
        raise HTTPException(status_code=400, detail="Empty file path")
    raw = _from_outputs_url(path_or_url)
    p = Path(raw)
    if not p.is_absolute():
        p = (get_project_root() / p).resolve()
    elif not p.exists():
        # 前端可能传了带 /outputs/ 的绝对形式，在服务端需按 project_root 解析（并解码 %40）
        raw_stripped = unquote(raw.lstrip("/"))
        if raw_stripped:
            p_rel = (get_project_root() / raw_stripped).resolve()
            if p_rel.exists():
                p = p_rel
    return p


def _convert_to_pdf(input_path: Path, output_dir: Path) -> Path:
    output_dir.mkdir(parents=True, exist_ok=True)
    cmd = [
        "libreoffice",
        "--headless",
        "--convert-to",
        "pdf",
        "--outdir",
        str(output_dir),
        str(input_path)
    ]
    subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
    pdf_path = output_dir / f"{input_path.stem}.pdf"
    if not pdf_path.exists():
        raise HTTPException(status_code=500, detail=f"PDF conversion failed for {input_path.name}")
    return pdf_path


def _merge_pdfs(pdf_paths: List[Path], output_path: Path) -> Path:
    if not pdf_paths:
        raise HTTPException(status_code=400, detail="No PDF files to merge")
    merged = fitz.open()
    for pdf in pdf_paths:
        with fitz.open(pdf) as src:
            merged.insert_pdf(src)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    merged.save(output_path)
    merged.close()
    return output_path


def _append_images_to_pptx(pptx_path: Path, image_paths: List[Path]) -> None:
    try:
        from pptx import Presentation
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"python-pptx not available: {e}")

    prs = Presentation(str(pptx_path))
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    for img_path in image_paths:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(img_path),
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height
        )
    prs.save(str(pptx_path))

@router.post("/upload")
async def upload_kb_file(
    file: UploadFile = File(...),
    email: str = Form(...),
    user_id: str = Form(...),
    notebook_id: Optional[str] = Form(None),
    notebook_title: Optional[str] = Form(None),
):
    """
    Upload a file to the notebook's knowledge base directory.
    New layout: outputs/{title}_{id}/sources/{stem}/original/
    Fallback: also writes to legacy kb_data path for backward compat.
    """
    if not email:
        raise HTTPException(status_code=400, detail="Email is required")
    if not notebook_id:
        raise HTTPException(status_code=400, detail="notebook_id is required for per-notebook storage")

    file_ext = Path(file.filename).suffix.lower()
    if file_ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(
            status_code=400,
            detail=f"Unsupported file type: {file_ext}. Allowed: {', '.join(ALLOWED_EXTENSIONS)}"
        )

    try:
        filename = file.filename or f"unnamed_{user_id}"
        filename = os.path.basename(filename)

        # --- New notebook-centric layout ---
        paths = get_notebook_paths(notebook_id, notebook_title or "", user_id)
        mgr = SourceManager(paths)

        # Save uploaded bytes to a temp location first, then import
        tmp_dir = paths.root / "_tmp"
        tmp_dir.mkdir(parents=True, exist_ok=True)
        tmp_path = tmp_dir / filename
        with open(tmp_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)

        source_info = await mgr.import_file(tmp_path, filename)

        # Clean up temp
        try:
            tmp_path.unlink(missing_ok=True)
        except Exception:
            pass

        # Build static URL from the original path in new layout
        project_root = get_project_root()
        rel = source_info.original_path.relative_to(project_root)
        static_path = "/" + rel.as_posix()

        # --- Also write to legacy path for backward compat ---
        legacy_dir = _notebook_dir(email, notebook_id)
        legacy_dir.mkdir(parents=True, exist_ok=True)
        legacy_path = legacy_dir / filename
        if not legacy_path.exists():
            shutil.copy2(str(source_info.original_path), str(legacy_path))

        # Auto-embed using new vector_store path
        embedded = False
        try:
            vector_base = str(paths.vector_store_dir)
            mineru_base = str(paths.source_mineru_dir(filename))
            file_list = [{"path": str(source_info.original_path)}]
            await process_knowledge_base_files(
                file_list=file_list,
                base_dir=vector_base,
                mineru_output_base=mineru_base,
            )
            embedded = True
            log.info("[upload] auto-embedding done: %s", filename)
        except Exception as emb_err:
            log.warning("[upload] auto-embedding failed for %s: %s", filename, emb_err)

        return {
            "success": True,
            "filename": filename,
            "file_size": os.path.getsize(source_info.original_path),
            "storage_path": str(source_info.original_path),
            "static_url": static_path,
            "file_type": file.content_type,
            "embedded": embedded,
        }

    except Exception as e:
        print(f"Error uploading file: {e}")
        raise HTTPException(status_code=500, detail=str(e))


def _sanitize_md_filename(title: str, prefix: str = "doc") -> str:
    """生成安全的 .md 文件名，避免路径注入与非法字符。"""
    safe = re.sub(r'[^\w\u4e00-\u9fff\s\-.]', "", (title or "").strip())
    safe = (safe or prefix)[:80].strip() or prefix
    return safe + f"_{int(time.time())}.md"


def _url_to_pdf(url: str, output_path: Path, timeout_ms: int = 30000) -> None:
    """
    使用 Playwright 打开 URL 并打印为 PDF，便于后续统一走 MinerU。
    若 Playwright 未安装或失败，抛出异常。
    """
    from playwright.sync_api import sync_playwright
    url = (url or "").strip()
    if not url.startswith(("http://", "https://")):
        raise ValueError("Invalid url")
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        try:
            page = browser.new_page()
            page.goto(url, wait_until="networkidle", timeout=timeout_ms)
            page.pdf(path=str(output_path), format="A4", print_background=True)
        finally:
            browser.close()


@router.post("/add-text-source")
async def add_text_source(
    notebook_id: str = Body(..., embed=True),
    email: str = Body(..., embed=True),
    user_id: Optional[str] = Body(None, embed=True),
    notebook_title: Optional[str] = Body(None, embed=True),
    title: str = Body("直接输入", embed=True),
    content: str = Body(..., embed=True),
) -> Dict[str, Any]:
    """
    将纯文本保存为笔记本内的 .md 文件并作为来源。用于「直接输入」引入。
    New layout: outputs/{title}_{id}/sources/{stem}/
    """
    if not notebook_id or not email:
        raise HTTPException(status_code=400, detail="notebook_id and email are required")
    if not (content or "").strip():
        raise HTTPException(status_code=400, detail="content is required")

    # New layout
    paths = get_notebook_paths(notebook_id, notebook_title or "", user_id)
    mgr = SourceManager(paths)
    source_info = await mgr.import_text(content, title)

    # Legacy compat
    user_dir = _notebook_dir(email, notebook_id)
    user_dir.mkdir(parents=True, exist_ok=True)
    legacy_filename = _sanitize_md_filename(title, "直接输入")
    legacy_path = user_dir / legacy_filename
    if not legacy_path.exists():
        try:
            shutil.copy2(str(source_info.original_path), str(legacy_path))
        except Exception:
            pass

    project_root = get_project_root()
    rel = source_info.original_path.relative_to(project_root)
    static_path = "/" + rel.as_posix()
    return {
        "success": True,
        "filename": source_info.original_path.name,
        "file_size": source_info.original_path.stat().st_size,
        "storage_path": str(source_info.original_path),
        "static_url": static_path,
        "id": f"file-{source_info.original_path.name}",
    }


@router.post("/import-url-as-source")
async def import_url_as_source(
    notebook_id: str = Body(..., embed=True),
    email: str = Body(..., embed=True),
    user_id: Optional[str] = Body(None, embed=True),
    notebook_title: Optional[str] = Body(None, embed=True),
    url: str = Body(..., embed=True),
) -> Dict[str, Any]:
    """
    抓取 URL 网页正文存为 .md 文件，作为来源。
    New layout: outputs/{title}_{id}/sources/{stem}/
    """
    if not notebook_id or not email:
        raise HTTPException(status_code=400, detail="notebook_id and email are required")
    url = (url or "").strip()
    if not url.startswith("http://") and not url.startswith("https://"):
        raise HTTPException(status_code=400, detail="Invalid url")

    # Fetch page text
    try:
        text = await asyncio.to_thread(fetch_page_text, url)
        if not text or text.startswith("[抓取失败"):
            raise RuntimeError(text or "fetch_page_text returned empty")
    except Exception as e:
        log.warning("fetch_page_text failed: %s", e)
        raise HTTPException(status_code=500, detail=f"网页抓取失败: {e}")

    # Parse title from URL
    try:
        parsed = urlparse(url)
        title = (parsed.netloc or "网页") + "_" + (parsed.path.strip("/") or "page")[:30]
    except Exception:
        title = "网页"

    # New layout
    paths = get_notebook_paths(notebook_id, notebook_title or "", user_id)
    mgr = SourceManager(paths)
    source_info = await mgr.import_url(url, text, title)

    # Legacy compat
    user_dir = _notebook_dir(email, notebook_id)
    user_dir.mkdir(parents=True, exist_ok=True)
    legacy_filename = _sanitize_md_filename(title, "网页")
    legacy_path = user_dir / legacy_filename
    if not legacy_path.exists():
        try:
            shutil.copy2(str(source_info.original_path), str(legacy_path))
        except Exception:
            pass

    project_root = get_project_root()
    rel = source_info.original_path.relative_to(project_root)
    static_path = "/" + rel.as_posix()
    return {
        "success": True,
        "filename": source_info.original_path.name,
        "file_size": source_info.original_path.stat().st_size,
        "storage_path": str(source_info.original_path),
        "static_url": static_path,
        "id": f"file-{source_info.original_path.name}",
    }


@router.delete("/delete")
async def delete_kb_file(
    storage_path: str = Form(...)
):
    """
    Delete a file from the physical storage.
    """
    try:
        # Security check: ensure path is within KB_BASE_DIR
        # This is a basic check. In production, use more robust path validation.
        target_path = Path(storage_path).resolve()
        base_path = KB_BASE_DIR.resolve()
        
        if not str(target_path).startswith(str(base_path)):
             # Allow if it's the absolute path provided by the user system
             # Check if it exists essentially
             pass

        if target_path.exists() and target_path.is_file():
            os.remove(target_path)
            return {"success": True, "message": "File deleted"}
        else:
            return {"success": False, "message": "File not found"}
            
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

def _vector_store_base_dir(email: Optional[str], notebook_id: Optional[str]) -> Optional[str]:
    """与 kb_embedding 约定一致：返回该 notebook 的向量库根目录，供 RAG 使用。"""
    root = get_project_root()
    if not email:
        base = root / "outputs" / "kb_data" / "vector_store_main"
    else:
        base = root / "outputs" / "kb_data" / (email or "default")
        if notebook_id:
            safe_nb = notebook_id.replace("/", "_").replace("\\", "_")[:128]
            base = base / safe_nb / "vector_store"
        else:
            base = base / "_shared" / "vector_store"
    return str(base) if base.exists() else None


@router.post("/chat")
async def chat_with_kb(
    files: List[str] = Body(..., embed=True),
    query: str = Body(..., embed=True),
    history: List[Dict[str, str]] = Body([], embed=True),
    email: Optional[str] = Body(None, embed=True),
    notebook_id: Optional[str] = Body(None, embed=True),
    api_url: Optional[str] = Body(None, embed=True),
    api_key: Optional[str] = Body(None, embed=True),
    model: str = Body(settings.KB_CHAT_MODEL, embed=True),
):
    """
    Intelligent QA Chat. 若传 email/notebook_id 且该 notebook 已建索引，会优先用 RAG 检索片段作为上下文。
    """
    try:
        # Normalize file paths (web path -> local absolute path)
        project_root = get_project_root()
        local_files = []
        for f in files:
            # remove leading /outputs/ if present, or just join
            # Web path: /outputs/kb_data/...
            clean_path = f.lstrip('/')
            p = project_root / clean_path
            if p.exists():
                local_files.append(str(p))
            else:
                # Try raw path
                p_raw = Path(f)
                if p_raw.exists():
                    local_files.append(str(p_raw))
        
        if not local_files:
             # Just return empty answer or handle logic
             pass

        vector_store_base_dir = _vector_store_base_dir(email, notebook_id)

        # Construct Request
        req = IntelligentQARequest(
            files=local_files,
            query=query,
            history=history,
            vector_store_base_dir=vector_store_base_dir,
            chat_api_url=api_url or os.getenv("DF_API_URL"),
            api_key=api_key or os.getenv("DF_API_KEY"),
            model=model
        )
        
        state = IntelligentQAState(request=req)
        
        # Run workflow via registry (统一使用 run_workflow)
        result_state = await run_workflow("intelligent_qa", state)
        
        # graph.ainvoke returns the final state dict or state object depending on implementation.
        # LangGraph usually returns dict. But our GenericGraphBuilder wrapper might return state.
        # GenericGraphBuilder compile returns a compiled graph.
        # Let's check typical usage. usually await graph.ainvoke(state) returns dict.
        
        answer = ""
        file_analyses = []
        source_mapping = {}

        if isinstance(result_state, dict):
            answer = result_state.get("answer", "")
            file_analyses = result_state.get("file_analyses", [])
            source_mapping = result_state.get("source_mapping", {})
        else:
            answer = getattr(result_state, "answer", "")
            file_analyses = getattr(result_state, "file_analyses", [])
            source_mapping = getattr(result_state, "source_mapping", {})

        # 将 source_mapping 的 int key 转为 str（JSON 要求）
        source_mapping_str = {str(k): v for k, v in source_mapping.items()} if source_mapping else {}

        return {
            "success": True,
            "answer": answer,
            "file_analyses": file_analyses,
            "source_mapping": source_mapping_str
        }

    except HTTPException:
        raise
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))


# ---------- 1.1 对话记录：入库与读取 ----------
def _supabase_upsert_conversation(email: str, user_id: Optional[str], notebook_id: Optional[str]) -> Optional[Dict[str, Any]]:
    sb = get_supabase_admin_client()
    if not sb:
        return None
    try:
        q = sb.table("kb_conversations").select("id,created_at,updated_at")
        if email:
            q = q.eq("user_email", email)
        if user_id:
            q = q.eq("user_id", user_id)
        if notebook_id:
            q = q.eq("notebook_id", notebook_id)
        else:
            q = q.is_("notebook_id", "null")
        r = q.order("updated_at", desc=True).limit(1).execute()
        rows = (r.data or []) if hasattr(r, "data") else []
        if rows:
            # update updated_at
            sb.table("kb_conversations").update({"updated_at": time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime())}).eq("id", rows[0]["id"]).execute()
            return rows[0]
        ins = sb.table("kb_conversations").insert({
            "user_email": email,
            "user_id": user_id,
            "notebook_id": notebook_id,
            "title": "对话",
        }).execute()
        data = (ins.data or []) if hasattr(ins, "data") else []
        return data[0] if data else None
    except Exception as e:
        log.warning("supabase conversation upsert failed: %s", e)
        return None


@router.get("/conversations")
async def list_conversations_get(
    email: Optional[str] = None,
    user_id: Optional[str] = None,
    notebook_id: Optional[str] = None,
) -> Dict[str, Any]:
    """List conversations for user (by email or user_id), optionally filter by notebook_id."""
    sb = get_supabase_admin_client()
    if not sb:
        return {"success": True, "conversations": []}
    try:
        q = sb.table("kb_conversations").select("id,notebook_id,title,created_at,updated_at")
        if email:
            q = q.eq("user_email", email)
        if user_id:
            q = q.eq("user_id", user_id)
        if notebook_id:
            q = q.eq("notebook_id", notebook_id)
        r = q.order("updated_at", desc=True).limit(50).execute()
        rows = (r.data or []) if hasattr(r, "data") else []
        return {"success": True, "conversations": rows}
    except Exception as e:
        log.warning("list_conversations failed: %s", e)
        return {"success": True, "conversations": []}


@router.post("/conversations")
async def get_or_create_conversation(
    email: str = Body(..., embed=True),
    user_id: Optional[str] = Body(None, embed=True),
    notebook_id: Optional[str] = Body(None, embed=True),
) -> Dict[str, Any]:
    """Get or create one conversation for this user+notebook. Returns conversation id."""
    conv = _supabase_upsert_conversation(email, user_id, notebook_id)
    if conv:
        return {"success": True, "conversation_id": conv.get("id"), "conversation": conv}
    return {"success": False, "conversation_id": None}


@router.get("/conversations/{conversation_id}/messages")
async def get_conversation_messages(conversation_id: str) -> Dict[str, Any]:
    sb = get_supabase_admin_client()
    if not sb:
        return {"success": True, "messages": []}
    try:
        r = sb.table("kb_chat_messages").select("id,role,content,created_at").eq(
            "conversation_id", conversation_id
        ).order("created_at", desc=False).execute()
        rows = (r.data or []) if hasattr(r, "data") else []
        return {"success": True, "messages": rows}
    except Exception as e:
        log.warning("get_conversation_messages failed: %s", e)
        return {"success": True, "messages": []}


@router.post("/conversations/{conversation_id}/messages")
async def append_conversation_messages(
    conversation_id: str,
    messages: List[Dict[str, str]] = Body(..., embed=True),
) -> Dict[str, Any]:
    """Append messages (list of {role, content})."""
    sb = get_supabase_admin_client()
    if not sb:
        return {"success": False, "message": "Database not configured"}
    try:
        rows = [{"conversation_id": conversation_id, "role": m.get("role", "user"), "content": m.get("content", "")} for m in messages]
        sb.table("kb_chat_messages").insert(rows).execute()
        return {"success": True}
    except Exception as e:
        log.warning("append_conversation_messages failed: %s", e)
        return {"success": False, "message": str(e)}


# ---------- 1.2 生成记录持久化：列表与写入 ----------
@router.get("/outputs")
async def list_outputs(
    email: Optional[str] = None,
    user_id: Optional[str] = None,
    notebook_id: Optional[str] = None,
    notebook_title: Optional[str] = None,
) -> Dict[str, Any]:
    """List generated outputs (ppt/mindmap/podcast/drawio) for user. Prefer DB, fallback to disk scan."""
    sb = get_supabase_admin_client()
    project_root = get_project_root()
    files: List[Dict[str, Any]] = []
    if sb:
        try:
            q = sb.table("kb_output_records").select("id,output_type,file_name,file_path,result_path,download_url,created_at")
            if email:
                q = q.eq("user_email", email)
            if user_id:
                q = q.eq("user_id", user_id)
            if notebook_id:
                q = q.eq("notebook_id", notebook_id)
            r = q.order("created_at", desc=True).limit(100).execute()
            rows = (r.data or []) if hasattr(r, "data") else []
            for row in rows:
                url = row.get("download_url") or row.get("result_path") or row.get("file_path")
                if url and not url.startswith("http") and not url.startswith("/"):
                    url = _to_outputs_url(url)
                files.append({
                    "id": row.get("id"),
                    "output_type": row.get("output_type", "ppt"),
                    "file_name": row.get("file_name"),
                    "download_url": url,
                    "created_at": row.get("created_at"),
                })
        except Exception as e:
            log.warning("list_outputs from db failed: %s", e)
    # Disk fallback: scan notebook-centric directory layout
    if not files and notebook_id:
        _FEATURE_EXT_MAP = {
            "ppt":     {".pdf", ".pptx"},
            "mindmap": {".mmd", ".mermaid"},
            "podcast": {".wav", ".mp3", ".m4a"},
            "drawio":  {".drawio"},
        }
        try:
            paths = get_notebook_paths(notebook_id, notebook_title or "", user_id)
            nb_root = paths.root
            if nb_root.exists():
                for feature, exts in _FEATURE_EXT_MAP.items():
                    feature_dir = nb_root / feature
                    if not feature_dir.exists():
                        continue
                    for ts_dir in feature_dir.iterdir():
                        if not ts_dir.is_dir():
                            continue
                        for f in ts_dir.iterdir():
                            if f.suffix.lower() in exts:
                                rel = str(f.relative_to(project_root))
                                files.append({
                                    "id": f"disk_{ts_dir.name}_{f.name}",
                                    "output_type": feature,
                                    "file_name": f.name,
                                    "download_url": _to_outputs_url(rel),
                                    "created_at": ts_dir.stat().st_mtime,
                                })
                                break
        except Exception as e:
            log.warning("list_outputs disk scan failed: %s", e)
    return {"success": True, "files": files}


def _extract_text_from_files(file_paths: List[str], max_chars: int = 50000) -> str:
    """从知识库文件列表中提取并合并文本，供 DrawIO 等使用。"""
    parts = []
    total = 0
    for f in file_paths:
        if total >= max_chars:
            break
        path = Path(f)
        if not path.exists():
            parts.append(f"[File not found: {f}]\n")
            total += len(parts[-1])
            continue
        suffix = path.suffix.lower()
        raw = ""
        try:
            if suffix == ".pdf":
                doc = fitz.open(f)
                raw = "\n".join(page.get_text() for page in doc)
                doc.close()
            elif suffix in [".docx", ".doc"]:
                try:
                    from docx import Document
                    doc = Document(f)
                    raw = "\n".join(p.text for p in doc.paragraphs)
                except Exception:
                    raw = "[Error: unsupported or missing python-docx]"
            elif suffix in [".pptx", ".ppt"]:
                try:
                    from pptx import Presentation
                    prs = Presentation(f)
                    raw = ""
                    for i, slide in enumerate(prs.slides):
                        raw += f"--- Slide {i+1} ---\n"
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                raw += shape.text + "\n"
                except Exception:
                    raw = "[Error: unsupported or missing python-pptx]"
            else:
                with open(path, "r", encoding="utf-8", errors="replace") as fp:
                    raw = fp.read()
        except Exception as e:
            raw = f"[Parse Error: {e}]"
        chunk = (raw[: max_chars - total] + ("..." if len(raw) > max_chars - total else "")) if raw else ""
        parts.append(chunk)
        total += len(chunk)
    return "\n\n".join(parts)


def _save_output_record(
    email: str,
    user_id: Optional[str],
    notebook_id: Optional[str],
    output_type: str,
    file_name: str,
    file_path: str,
    result_path: str,
    download_url: str,
):
    sb = get_supabase_admin_client()
    if not sb:
        return
    try:
        sb.table("kb_output_records").insert({
            "user_email": email,
            "user_id": user_id,
            "notebook_id": notebook_id,
            "output_type": output_type,
            "file_name": file_name,
            "file_path": file_path,
            "result_path": result_path,
            "download_url": download_url,
        }).execute()
    except Exception as e:
        log.warning("_save_output_record failed: %s", e)


# ---------- 1.3 笔记本（目录）与后端联动 ----------
def _notebooks_local_path(user_id: str) -> Path:
    root = get_project_root()
    base = root / "outputs" / "kb_data" / "_notebooks"
    base.mkdir(parents=True, exist_ok=True)
    safe_id = (user_id or "default").replace("/", "_").replace("\\", "_")[:64]
    return base / f"{safe_id}.json"


def _list_notebooks_local(user_id: str) -> List[Dict[str, Any]]:
    path = _notebooks_local_path(user_id)
    if not path.exists():
        return []
    try:
        raw = path.read_text(encoding="utf-8")
        data = json.loads(raw)
        return data if isinstance(data, list) else []
    except Exception as e:
        log.warning("_list_notebooks_local read failed: %s", e)
        return []


def _create_notebook_local(user_id: str, name: str, description: str = "") -> Dict[str, Any]:
    path = _notebooks_local_path(user_id)
    now = time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime())
    nb_id = f"local_{int(time.time() * 1000)}_{os.urandom(4).hex()}"
    new_nb = {
        "id": nb_id,
        "name": name,
        "description": description or "",
        "created_at": now,
        "updated_at": now,
    }
    notebooks = _list_notebooks_local(user_id)
    notebooks.insert(0, new_nb)
    try:
        path.write_text(json.dumps(notebooks, ensure_ascii=False, indent=2), encoding="utf-8")
    except Exception as e:
        log.warning("_create_notebook_local write failed: %s", e)
        raise HTTPException(status_code=500, detail="Failed to save notebook locally")
    return new_nb


# 不做用户管理时使用的默认用户，数据从 outputs 取
DEFAULT_USER_ID = "default"
DEFAULT_EMAIL = "default"


@router.get("/notebooks")
async def list_notebooks(
    email: Optional[str] = None,
    user_id: Optional[str] = None,
) -> Dict[str, Any]:
    """List notebooks. No user_id => use default (data from outputs)."""
    uid = (user_id or "").strip() or DEFAULT_USER_ID
    sb = get_supabase_admin_client()
    if sb:
        try:
            q = sb.table("knowledge_bases").select("id,name,description,created_at,updated_at").eq("user_id", uid)
            r = q.order("updated_at", desc=True).execute()
            rows = (r.data or []) if hasattr(r, "data") else []
            if rows:
                from collections import Counter
                nb_ids = [row["id"] for row in rows]
                try:
                    fr = sb.table("knowledge_base_files").select("kb_id").in_("kb_id", nb_ids).execute()
                    file_rows = (fr.data or []) if hasattr(fr, "data") else []
                    counts = Counter(f.get("kb_id") for f in file_rows if f.get("kb_id"))
                except Exception as e:
                    log.warning("notebooks file count failed: %s", e)
                    counts = {}
                for row in rows:
                    row["sources"] = counts.get(row["id"], 0)
            return {"success": True, "notebooks": rows}
        except Exception as e:
            log.warning("list_notebooks failed: %s", e)
            return {"success": True, "notebooks": []}
    rows = _list_notebooks_local(uid)
    email_for_path = (email or "").strip() or DEFAULT_EMAIL
    for row in rows:
        nb_id = row.get("id")
        if nb_id:
            count = 0
            # New layout count
            try:
                paths = get_notebook_paths(nb_id, row.get("name", ""), uid)
                if paths.sources_dir.exists():
                    count += sum(1 for d in paths.sources_dir.iterdir() if d.is_dir() and (d / "original").exists())
            except Exception:
                pass
            # Legacy count (avoid double-counting)
            nb_dir = _notebook_dir(email_for_path, nb_id)
            try:
                if nb_dir.exists():
                    file_count = sum(1 for _ in nb_dir.iterdir() if _.is_file() and _.name != LINK_SOURCES_FILENAME)
                    link_count = len(_load_link_sources(nb_dir))
                    count = max(count, file_count + link_count)
            except Exception:
                pass
            row["sources"] = count
        else:
            row.setdefault("sources", 0)
    return {"success": True, "notebooks": rows}


@router.get("/files")
async def list_notebook_files(
    user_id: Optional[str] = None,
    notebook_id: Optional[str] = None,
    email: Optional[str] = None,
    notebook_title: Optional[str] = None,
) -> Dict[str, Any]:
    """List files in a notebook. Reads from new sources/ dir first, then falls back to legacy kb_data/."""
    uid = (user_id or "").strip() or DEFAULT_USER_ID
    em = (email or "").strip() or DEFAULT_EMAIL
    if not notebook_id:
        return {"success": True, "files": []}

    files: List[Dict[str, Any]] = []
    seen_names: set = set()
    project_root = get_project_root()

    # --- 1) Read from new layout: outputs/{title}_{id}/sources/ ---
    try:
        paths = get_notebook_paths(notebook_id, notebook_title or "", uid)
        sources_dir = paths.sources_dir
        if sources_dir.exists():
            for src_dir in sorted(sources_dir.iterdir()):
                if not src_dir.is_dir():
                    continue
                orig_dir = src_dir / "original"
                if not orig_dir.exists():
                    continue
                for f in orig_dir.iterdir():
                    if not f.is_file():
                        continue
                    rel = f.relative_to(project_root)
                    static_url = "/" + rel.as_posix()
                    stat = f.stat()
                    files.append({
                        "id": f"file-{f.name}-{stat.st_mtime_ns}",
                        "name": f.name,
                        "url": static_url,
                        "static_url": static_url,
                        "file_size": stat.st_size,
                        "file_type": (f.suffix or "").lower() or "application/octet-stream",
                    })
                    seen_names.add(f.name)
    except Exception as e:
        log.warning("[list_notebook_files] new layout read failed: %s", e)

    # --- 2) Fallback: read from legacy kb_data/{email}/{notebook_id}/ ---
    nb_dir = _notebook_dir(em, notebook_id)
    link_static_urls: set = set()
    try:
        if nb_dir.exists():
            link_sources_path = nb_dir / LINK_SOURCES_FILENAME
            if link_sources_path.exists():
                try:
                    raw = link_sources_path.read_text(encoding="utf-8")
                    link_list = json.loads(raw)
                    if isinstance(link_list, list):
                        for item in link_list:
                            su = item.get("static_url") or ""
                            if su:
                                link_static_urls.add(su.rstrip("/"))
                except Exception:
                    pass
            for f in nb_dir.iterdir():
                if not f.is_file():
                    continue
                if f.name == LINK_SOURCES_FILENAME:
                    continue
                if f.name in seen_names:
                    continue
                rel = f.relative_to(project_root)
                static_url = "/" + rel.as_posix().replace("@", "%40")
                if static_url.rstrip("/") in link_static_urls:
                    continue
                stat = f.stat()
                files.append({
                    "id": f"file-{f.name}-{stat.st_mtime_ns}",
                    "name": f.name,
                    "url": static_url,
                    "static_url": static_url,
                    "file_size": stat.st_size,
                    "file_type": (f.suffix or "").lower() or "application/octet-stream",
                })
            # Link sources
            if link_sources_path.exists():
                try:
                    raw = link_sources_path.read_text(encoding="utf-8")
                    link_list = json.loads(raw)
                    if isinstance(link_list, list):
                        for i, item in enumerate(link_list):
                            link_id = item.get("id") or f"link-{i}-{hash(item.get('link', '')) % 10**8}"
                            static_url = item.get("static_url") or ""
                            link_url = item.get("link") or ""
                            url = static_url or link_url
                            files.append({
                                "id": link_id,
                                "name": (item.get("title") or item.get("link") or "Link")[:200],
                                "url": url,
                                "static_url": url,
                                "file_size": 0,
                                "file_type": "link",
                                "source_type": "link",
                                "snippet": item.get("snippet") or "",
                            })
                except Exception as e:
                    log.warning("list_notebook_files link_sources read failed: %s", e)
    except Exception as e:
        log.warning("list_notebook_files legacy read failed: %s", e)

    files.sort(key=lambda x: (x.get("file_type") == "link", x.get("name", "")))
    return {"success": True, "files": files}


@router.post("/fast-research")
async def fast_research(
    query: str = Body(..., embed=True),
    top_k: int = Body(10, embed=True),
    search_provider: Optional[str] = Body("serper", embed=True),
    search_api_key: Optional[str] = Body(None, embed=True),
    search_engine: Optional[str] = Body("google", embed=True),
    google_cse_id: Optional[str] = Body(None, embed=True),
) -> Dict[str, Any]:
    """
    Fast Research: 用户输入查询，搜索引擎搜索，返回 top_k 条结果作为候选来源。
    支持：
    - serper：Google（环境变量 SERPER_API_KEY）
    - serpapi：Google / 百度，传 search_api_key + search_engine（google | baidu）
    - google_cse：传 search_api_key + google_cse_id
    - brave：传 search_api_key
    - bocha：博查 AI 网页搜索（https://api.bocha.cn），传 search_api_key（Bearer 鉴权）
    """
    top_k = max(1, min(20, top_k))
    sources = fast_research_search(
        query,
        top_k=top_k,
        search_provider=search_provider or "serper",
        search_api_key=search_api_key,
        search_engine=search_engine or "google",
        google_cse_id=google_cse_id,
    )
    return {
        "success": True,
        "query": query,
        "sources": sources,
    }


def _pdf_to_markdown(local_path: str) -> str:
    """将本地 PDF 提取为可读文本/简单 markdown，用于前端展示。"""
    text_parts: List[str] = []
    try:
        doc = fitz.open(local_path)
        for i in range(len(doc)):
            page = doc[i]
            block_list = page.get_text("blocks")
            for block in block_list:
                # block: (x0, y0, x1, y1, "text", block_no, block_type)
                if len(block) >= 5 and block[4].strip():
                    text_parts.append(block[4].strip())
        doc.close()
    except Exception as e:
        log.warning("_pdf_to_markdown failed for %s: %s", local_path, e)
        return ""
    return "\n\n".join(text_parts)


def _manifest_path_for_storage_path(storage_path: str) -> Optional[Path]:
    """
    根据来源文件路径（如 /outputs/kb_data/default/notebook_id/xxx.pdf）推断该笔记本的
    vector_store 目录并返回 knowledge_manifest.json 路径；若无法推断则返回 None。
    """
    raw = _from_outputs_url((storage_path or "").strip())
    p = Path(raw)
    if not p.is_absolute():
        p = (get_project_root() / raw).resolve()
    parts = p.parts
    if "kb_data" not in parts:
        return None
    idx = parts.index("kb_data")
    if idx + 2 >= len(parts):
        return None
    email = parts[idx + 1]
    notebook_id = parts[idx + 2]
    root = get_project_root()
    safe_nb = (notebook_id or "_shared").replace("/", "_").replace("\\", "_")[:128]
    manifest_path = root / "outputs" / "kb_data" / email / safe_nb / "vector_store" / "knowledge_manifest.json"
    return manifest_path if manifest_path.exists() else None


@router.post("/get-source-display-content")
async def get_source_display_content(
    path: str = Body(..., embed=True),
    notebook_id: Optional[str] = Body(None, embed=True),
    email: Optional[str] = Body(None, embed=True),
) -> Dict[str, Any]:
    """
    返回用于前端展示的来源内容。若该来源已建索引且存在 MinerU 产出的 MD，则返回该 MD 内容；
    否则返回 from_mineru=false，前端可回退到 parse-local-file / fetch-page-content。
    """
    if not path or not path.strip():
        return {"content": None, "from_mineru": False}
    raw = _from_outputs_url(path.strip())
    abs_path = Path(raw)
    if not abs_path.is_absolute():
        abs_path = (get_project_root() / raw).resolve()
    if not abs_path.exists() or not abs_path.is_file():
        return {"content": None, "from_mineru": False}
    abs_str = str(abs_path.resolve())
    manifest_path = _manifest_path_for_storage_path(path)
    if not manifest_path:
        return {"content": None, "from_mineru": False}
    try:
        with open(manifest_path, "r", encoding="utf-8") as f:
            manifest = json.load(f)
    except Exception:
        return {"content": None, "from_mineru": False}
    for file_record in (manifest.get("files") or []):
        orig = (file_record.get("original_path") or "").strip()
        if not orig:
            continue
        if Path(orig).resolve() == abs_path.resolve():
            md_path = file_record.get("processed_md_path")
            if md_path and Path(md_path).exists():
                try:
                    content = Path(md_path).read_text(encoding="utf-8", errors="replace")
                    return {"content": content, "from_mineru": True}
                except Exception:
                    pass
            break
    return {"content": None, "from_mineru": False}


@router.post("/parse-local-file")
async def parse_local_file(path_or_url: str = Body(..., embed=True)) -> Dict[str, Any]:
    """
    解析本地文件内容，用于来源详情展示。支持 PDF：提取为文本/简单 markdown。
    path_or_url: 前端传来的路径，如 /outputs/kb_data/default/notebook_id/2025-6.pdf
    """
    if not path_or_url or not path_or_url.strip():
        raise HTTPException(status_code=400, detail="path_or_url is required")
    raw = _from_outputs_url(path_or_url.strip())
    p = Path(raw)
    if not p.is_absolute():
        p = (get_project_root() / raw).resolve()
    if not p.exists() or not p.is_file():
        raise HTTPException(status_code=404, detail="File not found")
    suffix = (p.suffix or "").lower()
    if suffix == ".pdf":
        content = _pdf_to_markdown(str(p))
        return {"success": True, "content": content or "[PDF 无文本或解析失败]", "format": "markdown"}
    if suffix in (".md", ".txt", ".markdown"):
        try:
            content = p.read_text(encoding="utf-8", errors="replace")
            return {"success": True, "content": content, "format": "markdown" if suffix == ".md" else "text"}
        except Exception as e:
            log.warning("parse_local_file read text failed: %s", e)
            raise HTTPException(status_code=500, detail=str(e))
    raise HTTPException(status_code=400, detail="Unsupported file type for preview (only .pdf, .md, .txt)")


@router.post("/fetch-page-content")
async def fetch_page_content(url: str = Body(..., embed=True)) -> Dict[str, Any]:
    """
    抓取 URL 对应页面的正文文本，用于来源详情展示（解析后的内容，可乱序但可读）。
    """
    if not url or not (url.startswith("http://") or url.startswith("https://")):
        raise HTTPException(status_code=400, detail="Invalid url")
    try:
        content = fetch_page_text(url, max_chars=50000)
        return {"success": True, "content": content}
    except Exception as e:
        log.warning("fetch_page_content failed: %s", e)
        raise HTTPException(status_code=500, detail=str(e))


def _load_link_sources(nb_dir: Path) -> List[Dict[str, Any]]:
    path = nb_dir / LINK_SOURCES_FILENAME
    if not path.exists():
        return []
    try:
        raw = path.read_text(encoding="utf-8")
        data = json.loads(raw)
        return data if isinstance(data, list) else []
    except Exception:
        return []


def _save_link_sources(nb_dir: Path, items: List[Dict[str, Any]]) -> None:
    nb_dir.mkdir(parents=True, exist_ok=True)
    path = nb_dir / LINK_SOURCES_FILENAME
    path.write_text(json.dumps(items, ensure_ascii=False, indent=2), encoding="utf-8")


def _resolve_link_to_local_md(email: str, notebook_id: Optional[str], link_url: str) -> Optional[Path]:
    """
    若该 link 在「引入」时已抓取并存为本地 .md，则返回该 .md 的 Path，否则返回 None。
    这样生成时优先用本地已存文件，不再重新爬。
    """
    if not link_url or not (link_url.startswith("http://") or link_url.startswith("https://")):
        return None
    nb_dir = _notebook_dir(email or "default", notebook_id)
    link_list = _load_link_sources(nb_dir)
    link_url_stripped = link_url.strip()
    for item in link_list:
        if (item.get("link") or "").strip() != link_url_stripped:
            continue
        static_url = (item.get("static_url") or "").strip()
        if not static_url:
            return None
        try:
            p = _resolve_local_path(static_url)
            if p.exists() and p.is_file():
                return p
        except Exception:
            pass
    return None


@router.post("/import-link-sources")
async def import_link_sources(
    notebook_id: str = Body(..., embed=True),
    email: str = Body(..., embed=True),
    user_id: Optional[str] = Body(None, embed=True),
    notebook_title: Optional[str] = Body(None, embed=True),
    items: List[Dict[str, Any]] = Body(..., embed=True),
) -> Dict[str, Any]:
    """
    将 Fast Research 等返回的候选来源导入到当前笔记本。
    每个 URL 通过 httpx 抓取正文存为 .md，然后自动触发 embedding。
    New layout: outputs/{title}_{id}/sources/{stem}/
    """
    if not notebook_id or not email:
        raise HTTPException(status_code=400, detail="notebook_id and email are required")

    # New layout
    paths = get_notebook_paths(notebook_id, notebook_title or "", user_id)
    mgr = SourceManager(paths)

    # Legacy compat
    nb_dir = _notebook_dir(email, notebook_id)
    nb_dir.mkdir(parents=True, exist_ok=True)
    existing = _load_link_sources(nb_dir)
    seen_links = {x.get("link") for x in existing if x.get("link")}
    imported = 0
    saved_md_paths: List[str] = []

    for it in items:
        link = (it.get("link") or "").strip()
        if not link or link in seen_links:
            continue
        title = (it.get("title") or "").strip() or link
        snippet = (it.get("snippet") or "").strip()
        static_url = ""
        filename = ""
        try:
            text = await asyncio.to_thread(fetch_page_text, link)
            if not text or text.startswith("[抓取失败"):
                raise RuntimeError(text or "empty response")

            # Import into new layout
            source_info = await mgr.import_url(link, text, title[:80])
            project_root = get_project_root()
            rel = source_info.original_path.relative_to(project_root)
            static_url = "/" + rel.as_posix()
            filename = source_info.original_path.name
            saved_md_paths.append(str(source_info.original_path))

            # Legacy compat: also save to old path
            legacy_filename = _sanitize_md_filename(title[:80], "link")
            legacy_path = nb_dir / legacy_filename
            if not legacy_path.exists():
                try:
                    shutil.copy2(str(source_info.original_path), str(legacy_path))
                except Exception:
                    pass

            log.info("[import-link-sources] 已抓取并保存: %s -> %s", link[:60], filename)
        except Exception as e:
            log.warning("[import-link-sources] 抓取失败 %s: %s", link[:60], e)

        existing.append({
            "id": f"link-{int(time.time() * 1000)}-{imported}",
            "title": title[:500],
            "link": link,
            "snippet": snippet[:2000],
            "static_url": static_url,
            "filename": filename,
        })
        seen_links.add(link)
        imported += 1
    _save_link_sources(nb_dir, existing)

    # Auto-embed saved .md files using new paths
    embedded = 0
    if saved_md_paths:
        try:
            vector_base = str(paths.vector_store_dir)
            file_list = [{"path": p} for p in saved_md_paths]
            await process_knowledge_base_files(
                file_list=file_list,
                base_dir=vector_base,
            )
            embedded = len(saved_md_paths)
            log.info("[import-link-sources] embedding 完成, %d 个文件", embedded)
        except Exception as e:
            log.warning("[import-link-sources] embedding 失败: %s", e)

    return {"success": True, "imported": imported, "embedded": embedded}


@router.post("/notebooks")
async def create_notebook(
    name: str = Body(..., embed=True),
    description: Optional[str] = Body(None, embed=True),
    user_id: str = Body(..., embed=True),
) -> Dict[str, Any]:
    """Create a notebook. Uses Supabase if configured, else local JSON file.
    Also creates the new outputs/{title}_{id}/sources/ directory."""
    sb = get_supabase_admin_client()
    nb_data = None
    if sb:
        try:
            ins = sb.table("knowledge_bases").insert({"user_id": user_id, "name": name, "description": description or ""}).execute()
            data = (ins.data or []) if hasattr(ins, "data") else []
            nb_data = data[0] if data else None
        except Exception as e:
            log.warning("create_notebook failed: %s", e)
            return {"success": False, "message": str(e)}
    else:
        try:
            nb_data = _create_notebook_local(user_id, name, description or "")
        except HTTPException:
            raise
        except Exception as e:
            log.warning("create_notebook local failed: %s", e)
            return {"success": False, "message": str(e)}

    # Create new directory structure
    if nb_data and nb_data.get("id"):
        try:
            paths = get_notebook_paths(nb_data["id"], name, user_id)
            paths.sources_dir.mkdir(parents=True, exist_ok=True)
            log.info("[create_notebook] created dir: %s", paths.root)
        except Exception as e:
            log.warning("[create_notebook] dir creation failed: %s", e)

    return {"success": True, "notebook": nb_data}


@router.post("/generate-ppt")
async def generate_ppt_from_kb(
    file_path: Optional[str] = Body(None, embed=True),
    file_paths: Optional[List[str]] = Body(None, embed=True),
    image_paths: Optional[List[str]] = Body(None, embed=True),
    image_items: Optional[List[Dict[str, Any]]] = Body(None, embed=True),
    query: Optional[str] = Body("", embed=True),
    need_embedding: bool = Body(False, embed=True),
    search_top_k: int = Body(8, embed=True),
    user_id: str = Body(..., embed=True),
    email: str = Body(..., embed=True),
    notebook_id: Optional[str] = Body(None, embed=True),
    notebook_title: Optional[str] = Body(None, embed=True),
    api_url: str = Body(..., embed=True),
    api_key: str = Body(..., embed=True),
    style: str = Body("modern", embed=True),
    language: str = Body("zh", embed=True),
    page_count: int = Body(10, embed=True),
    model: str = Body("deepseek-v3.2", embed=True),
    gen_fig_model: str = Body("gemini-2.5-flash-image", embed=True),
):
    """
    Generate PPT from knowledge base file. Outputs under user/notebook dir.
    """
    try:
        # 兼容前端传 file_paths 为数组或单个字符串；保证多选时每项一个来源
        if file_paths is not None:
            raw_list = file_paths if isinstance(file_paths, list) else [file_paths] if file_paths else []
        else:
            raw_list = [file_path] if file_path else []
        input_paths = [x for x in raw_list if x]
        log.info("[generate-ppt] 收到 file_paths 数量: %s", len(input_paths))

        if not input_paths:
            raise HTTPException(status_code=400, detail="No input files provided")
        if not isinstance(page_count, int) or page_count < 1 or page_count > 50:
            raise HTTPException(status_code=400, detail="page_count must be an integer between 1 and 50")
        log.info("[generate-ppt] 收到 page_count=%s", page_count)

        # 区分本地文件与网页 URL（前端会传 type=link 的 url 为 http(s) 链接）
        url_sources: List[str] = []
        path_sources: List[Path] = []
        user_image_items: List[Dict[str, Any]] = []
        seen_resolved: set = set()
        for p in input_paths:
            ps = (p or "").strip()
            if ps.startswith("http://") or ps.startswith("https://"):
                url_sources.append(ps)
                continue
            local_path = _resolve_local_path(p)
            if not local_path.exists():
                raise HTTPException(status_code=404, detail=f"File not found: {p}")
            ext = local_path.suffix.lower()
            if ext in IMAGE_EXTENSIONS:
                user_image_items.append({"path": str(local_path), "description": ""})
                continue
            if ext in {".pdf", ".pptx", ".ppt", ".docx", ".doc", ".md", ".markdown"}:
                key = str(local_path.resolve())
                if key not in seen_resolved:
                    seen_resolved.add(key)
                    path_sources.append(local_path)
            else:
                raise HTTPException(status_code=400, detail=f"Unsupported file type for PPT: {local_path.name}")

        if not path_sources and not url_sources:
            raise HTTPException(status_code=400, detail="At least one document or web source is required for PPT generation")

        ts = int(time.time())
        project_root = get_project_root()
        # New layout: outputs/{title}_{id}/ppt/{ts}/
        if notebook_id:
            nb_paths = get_notebook_paths(notebook_id, notebook_title or "", user_id)
            output_dir = nb_paths.feature_output_dir("ppt", ts)
        else:
            output_dir = _outputs_dir(email, notebook_id, f"{ts}_ppt")
        output_dir.mkdir(parents=True, exist_ok=True)

        md_exts = {".md", ".markdown"}
        pdf_like_exts = {".pdf", ".pptx", ".ppt", ".docx", ".doc"}
        doc_paths = path_sources
        md_paths = [p for p in doc_paths if p.suffix.lower() in md_exts]
        pdf_like_paths = [p for p in doc_paths if p.suffix.lower() in pdf_like_exts]

        use_text_input = bool(md_paths) or bool(url_sources)
        combined_text = ""
        local_file_path = None
        pdf_paths_for_outline: List[Path] = []

        if use_text_input:
            # 按 input_paths 顺序：先本地文件再 URL，生成「来源1」「来源2」…（含网页抓取）
            text_parts: List[str] = []
            idx = 0
            for p in input_paths:
                ps = (p or "").strip()
                if ps.startswith("http://") or ps.startswith("https://"):
                    idx += 1
                    content = None
                    local_md = _resolve_link_to_local_md(email, notebook_id, ps)
                    if local_md is not None:
                        try:
                            content = local_md.read_text(encoding="utf-8", errors="replace")
                            log.info("[generate-ppt] 网页来源使用已存 .md: %s", local_md.name)
                        except Exception as e:
                            log.warning("[generate-ppt] 读取已存 .md 失败 %s: %s", local_md, e)
                    if not (content or "").strip():
                        try:
                            content = fetch_page_text(ps, max_chars=100000)
                            if content:
                                log.info("[generate-ppt] 网页来源 %s 抓取成功，长度=%s", idx, len(content))
                        except Exception as e:
                            log.warning("[generate-ppt] 网页来源抓取失败 %s: %s", ps[:80], e)
                    if (content or "").strip():
                        text_parts.append(f"来源{len(text_parts) + 1}:\n{content.strip()}")
                    continue
                local_path = _resolve_local_path(p)
                if not local_path.exists():
                    continue
                ext = local_path.suffix.lower()
                if ext not in {".pdf", ".pptx", ".ppt", ".docx", ".doc", ".md", ".markdown"}:
                    continue
                idx += 1
                try:
                    if ext in md_exts:
                        content = local_path.read_text(encoding="utf-8")
                    elif ext == ".pdf":
                        # 优先从 MinerU 缓存读取高质量 markdown
                        content = _read_mineru_md_if_cached(local_path, email, notebook_id, notebook_title=notebook_title)
                        if not content:
                            content = _extract_text_from_files([str(local_path)])
                    elif ext in pdf_like_exts:
                        content = _extract_text_from_files([str(local_path)])
                    else:
                        content = ""
                    if (content or "").strip():
                        text_parts.append(f"来源{len(text_parts) + 1}:\n{content.strip()}")
                except Exception as e:
                    log.warning("read doc %s (来源%s): %s", local_path.name, len(text_parts) + 1, e)
            combined_text = "\n\n".join(text_parts).strip()
            log.info("[generate-ppt] 共 %s 个来源（本地 %s + 网页 %s），TEXT 块数: %s", len(path_sources) + len(url_sources), len(path_sources), len(url_sources), len(text_parts))
            if not combined_text:
                raise HTTPException(status_code=400, detail="No text content could be read from the selected sources")
        else:
            # 仅 PDF/PPTX/DOCX：转 PDF 后合并
            local_pdf_paths: List[Path] = []
            convert_dir = output_dir / "input"
            convert_dir.mkdir(parents=True, exist_ok=True)
            for p in pdf_like_paths:
                ext = p.suffix.lower()
                if ext == ".pdf":
                    local_pdf_paths.append(p)
                elif ext in {".pptx", ".ppt", ".docx", ".doc"}:
                    local_pdf_paths.append(_convert_to_pdf(p, convert_dir))
                else:
                    raise HTTPException(status_code=400, detail=f"Unsupported file type for PPT: {p.name}")

            pdf_paths_for_outline = local_pdf_paths
            if len(local_pdf_paths) > 1:
                merge_dir = output_dir / "input"
                merged_pdf = merge_dir / "merged.pdf"
                local_file_path = _merge_pdfs(local_pdf_paths, merged_pdf)
            else:
                local_file_path = local_pdf_paths[0]

        # Normalize image items (optional)
        resolved_image_items: List[Dict[str, Any]] = []
        for item in image_items or []:
            raw_path = item.get("path") or item.get("url") or ""
            if not raw_path:
                continue
            img_path = _resolve_local_path(str(raw_path))
            if img_path.exists() and img_path.suffix.lower() in IMAGE_EXTENSIONS:
                resolved_image_items.append({
                    "path": str(img_path),
                    "description": item.get("description") or item.get("desc") or ""
                })

        for img in image_paths or []:
            img_path = _resolve_local_path(img)
            if img_path.exists() and img_path.suffix.lower() in IMAGE_EXTENSIONS:
                resolved_image_items.append({
                    "path": str(img_path),
                    "description": ""
                })

        resolved_image_items.extend(user_image_items)

        # Embedding + retrieval (optional): use notebook-scoped vector store，入库只用本地 embedding，MinerU 输出到 kb_mineru
        retrieval_text = ""
        if need_embedding:
            base_dir = _notebook_dir(email, notebook_id) / "vector_store"
            embed_api_url = api_url
            if "/embeddings" not in embed_api_url:
                embed_api_url = embed_api_url.rstrip("/") + "/embeddings"
            project_root = get_project_root()
            safe_nb = (notebook_id or "_shared").replace("/", "_").replace("\\", "_")[:128]
            mineru_output_base = project_root / "outputs" / "kb_mineru" / (email or "default") / safe_nb
            mineru_output_base.mkdir(parents=True, exist_ok=True)

            files_for_embed = [{"path": str(p), "description": ""} for p in doc_paths]
            manifest = await process_knowledge_base_files(
                files_for_embed,
                base_dir=str(base_dir),
                api_url=None,
                api_key=api_key,
                model_name=None,
                multimodal_model=None,
                mineru_output_base=str(mineru_output_base),
            )

            manager = VectorStoreManager(
                base_dir=str(base_dir),
                api_key=api_key,
            )

            def _match_file_ids(m: Dict[str, Any], paths: List[Path]) -> List[str]:
                ids: List[str] = []
                target = {str(p.resolve()) for p in paths}
                for f in m.get("files", []):
                    try:
                        if str(Path(f.get("original_path", "")).resolve()) in target:
                            if f.get("id"):
                                ids.append(f["id"])
                    except Exception:
                        continue
                return ids

            file_ids = _match_file_ids(manifest or manager.manifest or {}, doc_paths)
            if query and file_ids:
                results = manager.search(query=query, top_k=search_top_k, file_ids=file_ids)
                retrieval_text = "\n\n".join([r.get("content", "") for r in results if r.get("content")])

        # Prepare request（支持 PDF 或 TEXT：.md 及混合时用 TEXT）
        ppt_req = Paper2PPTRequest(
            input_type="TEXT" if use_text_input else "PDF",
            input_content=combined_text if use_text_input else str(local_file_path),
            email=email,
            chat_api_url=api_url,
            chat_api_key=api_key,
            api_key=api_key,
            style=style,
            language=language,
            page_count=page_count,
            model=model,
            gen_fig_model=gen_fig_model,
            aspect_ratio="16:9",
            use_long_paper=False
        )
        log.info("[generate-ppt] ppt_req.page_count=%s（将传入 outline 生成）", ppt_req.page_count)

        # 复用 embedding 入库时已有的 MinerU 解析结果，避免重复跑 MinerU
        if not use_text_input and pdf_paths_for_outline:
            n_reused = _reuse_mineru_cache(pdf_paths_for_outline, output_dir, email, notebook_id, notebook_title=notebook_title)
            log.info("[generate-ppt] MinerU 缓存复用: %s/%s 个 PDF", n_reused, len(pdf_paths_for_outline))

        # Step 1: 生成大纲（kb_page_content 内含 LLM outline_agent，无人工确认）
        log.info("[generate-ppt] Step 1: 运行 kb_page_content，由 LLM 生成大纲 (outline)")
        state_pc = _init_state_from_request(ppt_req, result_path=output_dir)
        state_pc.kb_query = query or ""
        state_pc.kb_retrieval_text = retrieval_text
        state_pc.kb_user_images = resolved_image_items
        # 多 PDF 时按「来源1:\n...\n\n来源2:\n...」拼入，供 outline 使用
        if not use_text_input and len(pdf_paths_for_outline) > 1:
            multi_parts = []
            for i, p in enumerate(pdf_paths_for_outline):
                # 优先从 MinerU 缓存读取高质量 markdown
                part = _read_mineru_md_if_cached(p, email, notebook_id, notebook_title=notebook_title)
                if not part:
                    part = _extract_text_from_files([str(p)])
                if part.strip():
                    multi_parts.append(f"来源{i + 1}:\n{part}")
            if multi_parts:
                state_pc.kb_multi_source_text = "\n\n".join(multi_parts)
        state_pc_result = await run_workflow("kb_page_content", state_pc)
        if isinstance(state_pc_result, dict):
            for k, v in state_pc_result.items():
                setattr(state_pc, k, v)
        else:
            state_pc = state_pc_result
        pagecontent = getattr(state_pc, "pagecontent", []) or []
        log.info("[generate-ppt] Step 1 完成: 大纲已生成，共 %s 页", len(pagecontent))
        if not pagecontent:
            raise HTTPException(status_code=500, detail="大纲生成结果为空，请检查输入文档或重试")

        # Step 2: 按大纲生图并导出 PDF/PPTX（与 Paper2Any 一致使用 paper2ppt_parallel_consistent_style）
        state_pc.pagecontent = pagecontent
        log.info("[generate-ppt] Step 2: 运行 paper2ppt_parallel_consistent_style 生图")
        state_pp = await run_workflow("paper2ppt_parallel_consistent_style", state_pc)

        # Extract output paths (workflow may set ppt_pdf_path / ppt_pptx_path)
        if isinstance(state_pp, dict):
            pdf_path = state_pp.get("ppt_pdf_path") or ""
            pptx_path = state_pp.get("ppt_pptx_path") or ""
        else:
            pdf_path = getattr(state_pp, "ppt_pdf_path", None) or ""
            pptx_path = getattr(state_pp, "ppt_pptx_path", None) or ""
        # 若 workflow 未写回路径，则按约定路径回退：output_dir 下 paper2ppt.pdf / paper2ppt_editable.pptx
        if not pdf_path:
            fallback_pdf = Path(output_dir) / "paper2ppt.pdf"
            if fallback_pdf.exists():
                pdf_path = str(fallback_pdf)
        if not pptx_path:
            fallback_pptx = Path(output_dir) / "paper2ppt_editable.pptx"
            if fallback_pptx.exists():
                pptx_path = str(fallback_pptx)

        pdf_url = _to_outputs_url(pdf_path) if pdf_path else ""
        pptx_url = _to_outputs_url(pptx_path) if pptx_path else ""
        # 下载链接优先 PDF（可预览），其次 PPTX
        download_url = pdf_url or pptx_url
        _save_output_record(
            email=email,
            user_id=user_id,
            notebook_id=notebook_id,
            output_type="ppt",
            file_name="paper2ppt.pdf",
            file_path=pdf_path or "",
            result_path=str(output_dir),
            download_url=download_url,
        )

        return {
            "success": True,
            "result_path": str(output_dir),
            "pdf_path": pdf_url,
            "pptx_path": pptx_url,
            "download_url": download_url,
            "output_file_id": f"kb_ppt_{ts}",
        }

    except HTTPException:
        raise
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/generate-deep-research-report")
async def generate_deep_research_report(
    topic: str = Body(..., embed=True),
    user_id: str = Body(..., embed=True),
    email: str = Body(..., embed=True),
    notebook_id: Optional[str] = Body(None, embed=True),
    notebook_title: Optional[str] = Body(None, embed=True),
    api_url: str = Body(..., embed=True),
    api_key: str = Body(..., embed=True),
    language: str = Body("zh", embed=True),
    style: str = Body("modern", embed=True),
    page_count: int = Body(10, embed=True),
    model: str = Body("deepseek-v3.2", embed=True),
    gen_fig_model: str = Body("gemini-2.5-flash-image", embed=True),
    add_as_source: bool = Body(True, embed=True),
    search_provider: Optional[str] = Body("serper", embed=True),
    search_api_key: Optional[str] = Body(None, embed=True),
    search_engine: Optional[str] = Body("google", embed=True),
    search_top_k: int = Body(10, embed=True),
    # 新增：DeepResearch 完整模式配置
    use_full_deep_research: bool = Body(True, embed=True),  # 默认使用完整的阿里DeepResearch
    max_iterations: int = Body(50, embed=True),  # DeepResearch最大迭代次数
    serper_api_key: Optional[str] = Body(None, embed=True),  # Serper API密钥
    jina_api_key: Optional[str] = Body(None, embed=True),  # Jina API密钥
) -> Dict[str, Any]:
    """
    Deep Research 报告生成（默认使用完整版阿里DeepResearch）：
    - use_full_deep_research=True: 完整版（阿里DeepResearch多轮ReAct推理，深度）【默认】
    - use_full_deep_research=False: 简化版（搜索 + LLM总结，快速）
    """
    try:
        if not isinstance(page_count, int) or page_count < 1 or page_count > 50:
            raise HTTPException(status_code=400, detail="page_count must be an integer between 1 and 50")
        ts = int(time.time())
        project_root = get_project_root()

        # New layout: outputs/{title}_{id}/deep_research/{ts}/
        if notebook_id:
            dr_paths = get_notebook_paths(notebook_id, notebook_title or "", user_id)
            output_dir = dr_paths.feature_output_dir("deep_research", ts)
        else:
            output_dir = _outputs_dir(email, notebook_id, f"{ts}_deep_research")
        output_dir.mkdir(parents=True, exist_ok=True)

        topic = topic.strip()

        # ============================================================================
        # 模式选择：完整DeepResearch vs 简化版
        # ============================================================================

        if use_full_deep_research:
            # 使用完整的阿里DeepResearch（多轮ReAct推理）
            log.info("[generate-deep-research-report] 使用完整DeepResearch模式: topic=%r, max_iterations=%s", topic[:150], max_iterations)

            # 如果没有传递 serper_api_key，尝试使用 search_api_key 作为回退
            final_serper_key = serper_api_key or search_api_key

            log.info("[generate-deep-research-report] API配置: serper_api_key=%s, search_api_key=%s, final_serper_key=%s",
                     "***" if serper_api_key else "None",
                     "***" if search_api_key else "None",
                     "***" if final_serper_key else "None")

            # 运行完整DeepResearch（直接传递参数，不依赖环境变量）
            from fastapi_app.services.deep_research_integration import DeepResearchIntegration

            integration = DeepResearchIntegration(
                model_name=model,
                api_base=api_url,
                api_key=api_key,
                max_iterations=max_iterations,
                serper_key=final_serper_key,
                jina_keys=jina_api_key,
            )
            result = await integration.run_research(
                query=topic,
                max_iterations=max_iterations
            )

            if not result["success"]:
                raise HTTPException(status_code=500, detail=result.get("error", "DeepResearch failed"))

            # 格式化为Markdown
            report = integration.format_result_as_markdown(result)
            report_title = f"DeepResearch: {topic[:50]}"

            log.info("[generate-deep-research-report] 完整DeepResearch完成: iterations=%s, sources=%s",
                     result.get("iterations", 0), len(result.get("sources", [])))

        else:
            # 使用简化版（搜索 + LLM总结）
            search_top_k = max(1, min(20, search_top_k))
            log.info(
                "[generate-deep-research-report] 使用简化版模式: topic=%r, search_top_k=%s, provider=%s, model=%s",
                topic[:150], search_top_k, search_provider, model,
            )

            # 1) 搜索：用 topic 做 Fast Research，拿到 top_k 条结果
            sources = fast_research_search(
                topic,
                top_k=search_top_k,
                search_provider=search_provider or "serper",
                search_api_key=search_api_key or serper_api_key,
                search_engine=search_engine or "google",
            )
            log.info("[generate-deep-research-report] search 完成: 共 %s 条来源", len(sources))

            search_context = ""
            if sources:
                search_context = "\n\n".join(
                    f"[{i+1}] 标题: {s.get('title', '')}\n链接: {s.get('link', '')}\n摘要: {s.get('snippet', '')}"
                    for i, s in enumerate(sources)
                )
                log.info("[generate-deep-research-report] search_context 拼接完成: len=%s", len(search_context))
            else:
                log.warning("[generate-deep-research-report] no search results, LLM will generate from topic only")

            # 2) LLM：根据 topic + search_context 生成一篇长报告（返回标题 + 正文）
            report_title, report = generate_report_from_search(
                topic=topic,
                search_context=search_context,
                api_url=api_url,
                api_key=api_key,
                model=model,
                language=language,
            )
            if not (report or "").strip():
                raise HTTPException(status_code=500, detail="LLM did not return report content")
            log.info("[generate-deep-research-report] 简化版报告生成完成: title=%r, report_len=%s", report_title, len(report))

        # 3) 来源名：固定前缀 [report] + LLM 给的标题，保存为 .md
        safe_title = re.sub(r'[/\\:*?"<>|]', "", (report_title or "").strip()) or "report"
        safe_title = safe_title[:50].strip()
        file_name = f"[report] {safe_title}_{ts}.md"
        report_path = output_dir / file_name
        log.info("[generate-deep-research-report] 开始写入 Markdown: %s", report_path)
        report_path.write_text(report, encoding="utf-8")
        if not report_path.exists():
            raise HTTPException(status_code=500, detail="Deep research report file was not written")

        report_url = _to_outputs_url(str(report_path))
        log.info("[generate-deep-research-report] 报告已保存: %s, add_as_source=%s, notebook_id=%s", report_path, add_as_source, notebook_id)

        if add_as_source and notebook_id:
            nb_dir = _notebook_dir(email, notebook_id)
            nb_dir.mkdir(parents=True, exist_ok=True)
            dest = nb_dir / file_name
            shutil.copy2(str(report_path), dest)
            try:
                rel = dest.relative_to(project_root)
                source_static_url = "/" + rel.as_posix().replace("@", "%40")
            except ValueError:
                source_static_url = report_url
            _save_output_record(
                email=email,
                user_id=user_id,
                notebook_id=notebook_id,
                output_type="report",
                file_name=file_name,
                file_path=str(dest),
                result_path=str(output_dir),
                download_url=report_url,
            )
            stat = dest.stat()
            added_file = {
                "id": f"file-{file_name}-{stat.st_mtime_ns}",
                "name": file_name,
                "url": source_static_url,
                "static_url": source_static_url,
                "file_size": stat.st_size,
                "file_type": "text/markdown",
            }
            log.info("[generate-deep-research-report] 完成: 已加入来源, file_name=%s", file_name)
            return {
                "success": True,
                "pdf_path": report_url,
                "pdf_url": report_url,
                "file_name": file_name,
                "source_static_url": source_static_url,
                "added_as_source": True,
                "added_file": added_file,
            }

        _save_output_record(
            email=email,
            user_id=user_id,
            notebook_id=notebook_id,
            output_type="report",
            file_name=file_name,
            file_path=str(report_path),
            result_path=str(output_dir),
            download_url=report_url,
        )
        log.info("[generate-deep-research-report] 完成: 未加入来源, file_name=%s", file_name)
        return {
            "success": True,
            "pdf_path": report_url,
            "pdf_url": report_url,
            "file_name": file_name,
            "added_as_source": False,
        }
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/generate-podcast")
async def generate_podcast_from_kb(
    file_paths: List[str] = Body(..., embed=True),
    user_id: str = Body(..., embed=True),
    email: str = Body(..., embed=True),
    notebook_id: Optional[str] = Body(None, embed=True),
    notebook_title: Optional[str] = Body(None, embed=True),
    api_url: str = Body(..., embed=True),
    api_key: str = Body(..., embed=True),
    model: str = Body("deepseek-v3.2", embed=True),
    tts_model: str = Body("gemini-2.5-pro-preview-tts", embed=True),
    voice_name: str = Body("Kore", embed=True),
    voice_name_b: str = Body("Puck", embed=True),
    podcast_mode: str = Body("monologue", embed=True),
    language: str = Body("zh", embed=True),
):
    """
    从知识库生成播客。支持本地文件与「搜索引入」的 URL：URL 优先用已存 .md，否则抓取后写临时 .md 再参与生成。
    """
    try:
        ts = int(time.time())
        # New layout: outputs/{title}_{id}/podcast/{ts}/
        if notebook_id:
            paths = get_notebook_paths(notebook_id, notebook_title or "", user_id)
            output_dir = paths.feature_output_dir("podcast", ts)
        else:
            output_dir = _outputs_dir(email, notebook_id, f"{ts}_podcast")
        output_dir.mkdir(parents=True, exist_ok=True)
        project_root = get_project_root()

        if not file_paths:
            raise HTTPException(status_code=400, detail="No valid files provided")

        local_paths: List[Path] = []
        for f in (file_paths or []):
            ps = (f or "").strip()
            if ps.startswith("http://") or ps.startswith("https://"):
                content = None
                local_md = _resolve_link_to_local_md(email, notebook_id, ps)
                if local_md is not None:
                    try:
                        content = local_md.read_text(encoding="utf-8", errors="replace")
                        log.info("[generate-podcast] 网页来源使用已存 .md: %s", local_md.name)
                    except Exception as e:
                        log.warning("[generate-podcast] 读取已存 .md 失败: %s", e)
                if not (content or "").strip():
                    try:
                        content = fetch_page_text(ps, max_chars=100000)
                    except Exception as e:
                        log.warning("[generate-podcast] 抓取 URL 失败 %s: %s", ps[:60], e)
                        content = ""
                if (content or "").strip():
                    link_dir = output_dir / "input"
                    link_dir.mkdir(parents=True, exist_ok=True)
                    tmp_md = link_dir / f"link_{len(local_paths)}.md"
                    tmp_md.write_text(content.strip(), encoding="utf-8")
                    local_paths.append(tmp_md)
            else:
                local_path = _resolve_local_path(ps)
                if not local_path.exists() or not local_path.is_file():
                    raise HTTPException(status_code=404, detail=f"File not found: {ps}")
                local_paths.append(local_path)

        # 过滤不支持的文件类型（例如图片），只保留可转文本的文档（含 .md 报告）
        supported_exts = {".pdf", ".docx", ".doc", ".pptx", ".ppt", ".md", ".markdown"}
        filtered_paths: List[Path] = []
        ignored_paths: List[Path] = []
        for p in local_paths:
            if p.suffix.lower() in supported_exts:
                filtered_paths.append(p)
            else:
                ignored_paths.append(p)

        if not filtered_paths:
            raise HTTPException(status_code=400, detail="No supported document files for podcast (support: PDF, Word, PPT, MD)")

        if ignored_paths:
            log.warning(
                "[kb_podcast] ignore unsupported files: "
                + ", ".join([p.name for p in ignored_paths])
            )

        # If multiple files, merge into a single PDF (doc/ppt will be converted); 仅 .md 时拼成一份 .md
        if len(filtered_paths) > 1:
            merge_dir = output_dir / "input"
            merge_dir.mkdir(parents=True, exist_ok=True)

            pdf_paths: List[Path] = []
            md_paths: List[Path] = []
            for p in filtered_paths:
                ext = p.suffix.lower()
                if ext == ".pdf":
                    pdf_paths.append(p)
                elif ext in {".docx", ".doc", ".pptx", ".ppt"}:
                    pdf_paths.append(_convert_to_pdf(p, merge_dir))
                elif ext in {".md", ".markdown"}:
                    md_paths.append(p)

            if pdf_paths:
                merged_pdf = merge_dir / "merged.pdf"
                local_file_paths = [str(_merge_pdfs(pdf_paths, merged_pdf))]
            elif md_paths:
                merged_md = merge_dir / "merged.md"
                parts = [p.read_text(encoding="utf-8", errors="replace") for p in md_paths]
                merged_md.write_text("\n\n".join(parts), encoding="utf-8")
                local_file_paths = [str(merged_md)]
            else:
                raise HTTPException(status_code=400, detail="No supported document files for podcast")
        else:
            local_file_paths = [str(filtered_paths[0])]

        # Prepare request
        podcast_req = KBPodcastRequest(
            files=local_file_paths,
            chat_api_url=api_url,
            api_key=api_key,
            model=model,
            tts_model=tts_model,
            voice_name=voice_name,
            voice_name_b=voice_name_b,
            podcast_mode=podcast_mode,
            language=language
        )
        podcast_req.email = email

        state = KBPodcastState(request=podcast_req, result_path=str(output_dir))

        # Run workflow via registry (统一使用 run_workflow)
        result_state = await run_workflow("kb_podcast", state)

        # Extract results
        audio_path = ""
        script_path = ""
        result_path = ""

        if isinstance(result_state, dict):
            audio_path = result_state.get("audio_path", "")
            result_path = result_state.get("result_path", "")
        else:
            audio_path = getattr(result_state, "audio_path", "")
            result_path = getattr(result_state, "result_path", "")

        if result_path:
            script_path = str(Path(result_path) / "script.txt")

        audio_error = ""
        if not audio_path:
            audio_error = "No audio path returned from workflow"
        elif isinstance(audio_path, str) and audio_path.startswith("["):
            audio_error = audio_path
        else:
            audio_file = Path(audio_path)
            if not audio_file.is_absolute():
                audio_file = (get_project_root() / audio_file).resolve()
            if not audio_file.exists():
                audio_error = f"Audio file not found: {audio_file}"

        if audio_error:
            raise HTTPException(status_code=500, detail=audio_error)

        audio_url = _to_outputs_url(audio_path) if audio_path else ""
        script_url = _to_outputs_url(script_path) if script_path else ""
        result_url = _to_outputs_url(result_path) if result_path else ""

        _save_output_record(
            email=email,
            user_id=user_id,
            notebook_id=notebook_id,
            output_type="podcast",
            file_name=Path(audio_path).name if audio_path else "podcast.wav",
            file_path=audio_path or "",
            result_path=result_path or str(output_dir),
            download_url=audio_url,
        )

        return {
            "success": True,
            "result_path": result_url,
            "audio_path": audio_url,
            "script_path": script_url,
            "output_file_id": f"kb_podcast_{int(time.time())}"
        }

    except HTTPException:
        raise
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

@router.post("/generate-mindmap")
async def generate_mindmap_from_kb(
    file_paths: List[str] = Body(..., embed=True),
    user_id: str = Body(..., embed=True),
    email: str = Body(..., embed=True),
    notebook_id: Optional[str] = Body(None, embed=True),
    notebook_title: Optional[str] = Body(None, embed=True),
    api_url: str = Body(..., embed=True),
    api_key: str = Body(..., embed=True),
    model: str = Body("deepseek-v3.2", embed=True),
    mindmap_style: str = Body("default", embed=True),
    max_depth: int = Body(3, embed=True),
    language: str = Body("zh", embed=True),
):
    """
    从知识库生成思维导图。支持本地文件与「搜索引入」的 URL：路径用 _resolve_local_path；URL 优先用已存 .md，否则抓取后写临时 .md。
    """
    try:
        project_root = get_project_root()
        ts = int(time.time())
        # New layout: outputs/{title}_{id}/mindmap/{ts}/
        if notebook_id:
            paths = get_notebook_paths(notebook_id, notebook_title or "", user_id)
            output_dir = paths.feature_output_dir("mindmap", ts)
        else:
            output_dir = _outputs_dir(email, notebook_id, f"{ts}_mindmap_input")
        output_dir.mkdir(parents=True, exist_ok=True)
        local_file_paths: List[str] = []

        for f in (file_paths or []):
            ps = (f or "").strip()
            if ps.startswith("http://") or ps.startswith("https://"):
                content = None
                local_md = _resolve_link_to_local_md(email, notebook_id, ps)
                if local_md is not None:
                    try:
                        content = local_md.read_text(encoding="utf-8", errors="replace")
                        log.info("[generate-mindmap] 网页来源使用已存 .md: %s", local_md.name)
                    except Exception as e:
                        log.warning("[generate-mindmap] 读取已存 .md 失败: %s", e)
                if not (content or "").strip():
                    try:
                        content = fetch_page_text(ps, max_chars=100000)
                    except Exception as e:
                        log.warning("[generate-mindmap] 抓取 URL 失败 %s: %s", ps[:60], e)
                        content = ""
                if (content or "").strip():
                    tmp_md = output_dir / f"link_{len(local_file_paths)}.md"
                    tmp_md.write_text(content.strip(), encoding="utf-8")
                    local_file_paths.append(str(tmp_md))
            else:
                local_path = _resolve_local_path(ps)
                if not local_path.exists() or not local_path.is_file():
                    raise HTTPException(status_code=404, detail=f"File not found: {ps}")
                local_file_paths.append(str(local_path))

        if not local_file_paths:
            raise HTTPException(status_code=400, detail="No valid files provided")

        # Prepare request
        mindmap_req = KBMindMapRequest(
            files=local_file_paths,
            chat_api_url=api_url,
            api_key=api_key,
            model=model,
            mindmap_style=mindmap_style,
            max_depth=max_depth,
            language=language
        )
        mindmap_req.email = email

        state = KBMindMapState(request=mindmap_req)

        # Run workflow via registry (统一使用 run_workflow)
        result_state = await run_workflow("kb_mindmap", state)

        # Extract results
        mermaid_code = ""
        result_path = ""

        if isinstance(result_state, dict):
            mermaid_code = result_state.get("mermaid_code", "")
            result_path = result_state.get("result_path", "")
        else:
            mermaid_code = getattr(result_state, "mermaid_code", "")
            result_path = getattr(result_state, "result_path", "")

        mindmap_path = ""
        if result_path:
            mmd_path = Path(result_path) / "mindmap.mmd"
            if (not mmd_path.exists()) and mermaid_code:
                try:
                    mmd_path.write_text(mermaid_code, encoding="utf-8")
                except Exception:
                    pass
            if mmd_path.exists():
                mindmap_path = _to_outputs_url(str(mmd_path))

        _save_output_record(
            email=email,
            user_id=user_id,
            notebook_id=notebook_id,
            output_type="mindmap",
            file_name="mindmap.mmd",
            file_path=str(Path(result_path) / "mindmap.mmd") if result_path else "",
            result_path=result_path or "",
            download_url=mindmap_path,
        )

        return {
            "success": True,
            "result_path": _to_outputs_url(result_path) if result_path else "",
            "mermaid_code": mermaid_code,
            "mindmap_path": mindmap_path,
            "output_file_id": f"kb_mindmap_{int(time.time())}"
        }

    except HTTPException:
        raise
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))


def _collect_figure_images(
    mgr: "SourceManager",
    file_paths: List[str],
    project_root: Path,
) -> List[tuple]:
    """
    从 file_paths 对应的 source 中收集 MinerU 提取的 figure 图片。
    返回 [(source_stem, image_path), ...]
    """
    IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".bmp", ".tiff", ".webp"}
    results: List[tuple] = []

    for fp in (file_paths or []):
        ps = (fp or "").strip()
        if ps.startswith("http://") or ps.startswith("https://"):
            continue

        local_path = _resolve_local_path(ps)
        stem = local_path.stem

        # 1) 检查 MinerU auto/images/ 目录下的图片
        mineru_root = mgr.get_mineru_root(stem)
        if mineru_root and mineru_root.exists():
            images_dir = mineru_root / "images"
            scan_dir = images_dir if images_dir.is_dir() else mineru_root
            for img in sorted(scan_dir.iterdir()):
                if img.is_file() and img.suffix.lower() in IMAGE_EXTS:
                    results.append((stem, img))

        # 2) 如果 MinerU 没有图片，检查原始文件本身是否是图片
        if not any(s == stem for s, _ in results):
            if local_path.exists() and local_path.suffix.lower() in IMAGE_EXTS:
                results.append((stem, local_path))

    return results


@router.post("/generate-drawio")
async def generate_drawio_from_kb(
    file_paths: List[str] = Body(..., embed=True),
    user_id: str = Body(..., embed=True),
    email: str = Body(..., embed=True),
    notebook_id: Optional[str] = Body(None, embed=True),
    notebook_title: Optional[str] = Body(None, embed=True),
    api_url: str = Body(..., embed=True),
    api_key: str = Body(..., embed=True),
    model: str = Body("deepseek-v3.2", embed=True),
    diagram_type: str = Body("auto", embed=True),
    diagram_style: str = Body("default", embed=True),
    language: str = Body("zh", embed=True),
    text_content: Optional[str] = Body(None, embed=True),
):
    """
    从知识库选中文件生成 DrawIO 图表。
    优先从 MinerU 提取的 figure 图片走 SAM3 分割生成 drawio（缓存到 sources/{stem}/sam3/），
    没有 figure 图片时 fallback 到文本模式 LLM 生成。
    """
    try:
        log.info("[generate-drawio] 收到 file_paths: %s", file_paths)
        project_root = get_project_root()

        # --- SAM3 图片模式：有 notebook 且能找到 figure 图片时自动走 SAM3 ---
        if notebook_id:
            from fastapi_app.services.paper2drawio_service import Paper2DrawioService

            paths = get_notebook_paths(notebook_id, notebook_title or "", email)
            mgr = SourceManager(paths)
            ts = int(time.time())
            output_dir = paths.feature_output_dir("drawio", ts)
            output_dir.mkdir(parents=True, exist_ok=True)

            # 收集所有 source 中的 figure 图片
            figure_images = _collect_figure_images(mgr, file_paths, project_root)
            if not figure_images:
                log.warning("[generate-drawio] SAM3 模式但未找到 figure 图片，回退到文本模式")
            else:
                service = Paper2DrawioService()
                all_xmls = []
                for stem, img_path in figure_images:
                    sam3_cache = str(mgr.ensure_sam3_dir(stem))
                    result = await service.generate_diagram_from_image(
                        image_path=str(img_path),
                        chat_api_url=api_url,
                        api_key=api_key,
                        model=model,
                        language=language,
                        email=email,
                        sam3_cache_dir=sam3_cache,
                        output_dir=str(output_dir),
                    )
                    if result.get("success") and result.get("xml_content"):
                        all_xmls.append(result["xml_content"])
                        log.info("[generate-drawio] SAM3 成功: stem=%s", stem)
                    else:
                        log.warning("[generate-drawio] SAM3 失败: stem=%s err=%s", stem, result.get("error"))

                if all_xmls:
                    xml_content = all_xmls[0]  # 目前取第一个成功的
                    drawio_path = output_dir / f"diagram_{ts}.drawio"
                    drawio_path.write_text(xml_content, encoding="utf-8")
                    download_url = _to_outputs_url(str(drawio_path))
                    _save_output_record(
                        email=email, user_id=user_id, notebook_id=notebook_id,
                        output_type="drawio", file_name=drawio_path.name,
                        file_path=str(drawio_path), result_path=str(output_dir),
                        download_url=download_url,
                    )
                    return {
                        "success": True,
                        "xml_content": xml_content,
                        "file_path": download_url,
                        "error": None,
                        "output_file_id": f"kb_drawio_{ts}",
                    }
                # SAM3 全部失败，fall through 到文本模式
                log.warning("[generate-drawio] SAM3 全部失败，回退到文本模式")

        # --- 文本模式（原有逻辑） ---
        url_sources = []
        local_file_paths = []
        for f in (file_paths or []):
            ps = (f or "").strip()
            if ps.startswith("http://") or ps.startswith("https://"):
                url_sources.append(ps)
            else:
                local_path = _resolve_local_path(ps)
                if not local_path.exists() or not local_path.is_file():
                    log.warning("[generate-drawio] 文件不存在: 原始=%s 解析后=%s", ps, local_path)
                    raise HTTPException(status_code=404, detail=f"File not found: {ps}")
                local_file_paths.append(str(local_path))

        parts = []
        for i, url in enumerate(url_sources):
            # 优先用引入时已存的 .md，不重新爬
            local_md = _resolve_link_to_local_md(email, notebook_id, url)
            if local_md is not None:
                try:
                    content = local_md.read_text(encoding="utf-8", errors="replace")
                    if content.strip():
                        parts.append(f"来源{i + 1}:\n{content}")
                        log.info("[generate-drawio] 使用已存 .md: %s", local_md.name)
                        continue
                except Exception as e:
                    log.warning("[generate-drawio] 读取已存 .md 失败 %s: %s", local_md, e)
            try:
                content = fetch_page_text(url, max_chars=100000)
                if content and not content.startswith("["):
                    parts.append(f"来源{i + 1}:\n{content}")
                else:
                    parts.append(f"来源{i + 1}:\n[抓取失败或无正文]")
            except Exception as e:
                log.warning("[generate-drawio] 抓取 URL 失败 %s: %s", url[:60], e)
                parts.append(f"来源{i + 1}:\n[抓取失败: {e}]")
        if local_file_paths:
            local_text = _extract_text_from_files(local_file_paths)
            if local_text.strip():
                parts.append(local_text)
        text_content = "\n\n".join(parts) if parts else ""
        if not text_content.strip():
            raise HTTPException(
                status_code=400,
                detail="No text from selected sources (URL fetch failed or files empty). Check link or choose local files.",
            )

        from fastapi_app.services.paper2drawio_service import Paper2DrawioService

        service = Paper2DrawioService()
        result = await service.generate_diagram(
            request=None,
            chat_api_url=api_url,
            api_key=api_key,
            model=model,
            enable_vlm_validation=False,
            vlm_model=getattr(settings, "PAPER2DRAWIO_VLM_MODEL", "deepseek-v3.2"),
            vlm_validation_max_retries=3,
            input_type="TEXT",
            diagram_type=diagram_type,
            diagram_style=diagram_style,
            language=language,
            email=email,
            file=None,
            text_content=text_content,
        )

        if not result.get("success") or not result.get("xml_content"):
            return {
                "success": False,
                "xml_content": "",
                "file_path": "",
                "error": result.get("error") or "Failed to generate diagram",
                "output_file_id": None,
            }

        xml_content = result["xml_content"]
        ts = int(time.time())
        # New layout: outputs/{title}_{id}/drawio/{ts}/
        if notebook_id:
            paths = get_notebook_paths(notebook_id, notebook_title or "", user_id)
            output_dir = paths.feature_output_dir("drawio", ts)
        else:
            output_dir = project_root / OUTPUTS_BASE / (email or "default") / "_shared" / "drawio"
        output_dir.mkdir(parents=True, exist_ok=True)
        drawio_path = output_dir / f"diagram_{ts}.drawio"
        drawio_path.write_text(xml_content, encoding="utf-8")
        download_url = _to_outputs_url(str(drawio_path))

        _save_output_record(
            email=email,
            user_id=user_id,
            notebook_id=notebook_id,
            output_type="drawio",
            file_name=drawio_path.name,
            file_path=str(drawio_path),
            result_path=str(output_dir),
            download_url=download_url,
        )

        return {
            "success": True,
            "xml_content": xml_content,
            "file_path": download_url,
            "error": None,
            "output_file_id": f"kb_drawio_{ts}",
        }
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/save-mindmap")
async def save_mindmap_to_file(
    file_url: str = Body(..., embed=True),
    content: str = Body(..., embed=True),
):
    """
    Save edited Mermaid mindmap code back to the output file.
    """
    try:
        if not file_url:
            raise HTTPException(status_code=400, detail="File URL is required")

        local_path = Path(_from_outputs_url(file_url))
        if not local_path.is_absolute():
            local_path = (get_project_root() / local_path).resolve()

        outputs_root = (get_project_root() / "outputs").resolve()
        try:
            local_path.relative_to(outputs_root)
        except ValueError:
            raise HTTPException(status_code=400, detail="Invalid output path")

        if local_path.suffix.lower() not in {".mmd", ".mermaid", ".md"}:
            raise HTTPException(status_code=400, detail="Invalid mindmap file type")

        local_path.parent.mkdir(parents=True, exist_ok=True)
        local_path.write_text(content or "", encoding="utf-8")

        return {
            "success": True,
            "mindmap_path": _to_outputs_url(str(local_path))
        }
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))


# ===================== Flashcard 闪卡 =====================

@router.post("/generate-flashcards")
async def generate_flashcards(
    file_paths: List[str] = Body(..., embed=True),
    email: str = Body(..., embed=True),
    user_id: str = Body(..., embed=True),
    notebook_id: Optional[str] = Body(None, embed=True),
    notebook_title: Optional[str] = Body(None, embed=True),
    api_url: str = Body(..., embed=True),
    api_key: str = Body(..., embed=True),
    model: str = Body("deepseek-v3.2", embed=True),
    language: str = Body("zh", embed=True),
    card_count: int = Body(20, embed=True),
):
    """从知识库文件生成闪卡"""
    try:
        from fastapi_app.services.flashcard_service import generate_flashcards_with_llm

        local_paths = []
        for f in file_paths:
            ps = (f or "").strip()
            if ps.startswith("http://") or ps.startswith("https://"):
                local_md = _resolve_link_to_local_md(email, notebook_id, ps)
                if local_md and local_md.exists():
                    local_paths.append(str(local_md))
            else:
                local_path = _resolve_local_path(f)
                if local_path.exists():
                    local_paths.append(str(local_path))

        if not local_paths:
            raise HTTPException(status_code=400, detail="No valid files provided")

        text_content = _extract_text_from_files(local_paths, max_chars=50000)
        if not text_content.strip():
            raise HTTPException(status_code=400, detail="No text content extracted")

        log.info("[generate-flashcards] text_len=%d, files=%d", len(text_content), len(local_paths))

        flashcards = await generate_flashcards_with_llm(
            text_content=text_content,
            api_url=api_url,
            api_key=api_key,
            model=model,
            language=language,
            card_count=card_count,
        )
        if not flashcards:
            raise HTTPException(status_code=500, detail="Failed to generate flashcards")

        ts = int(time.time())
        flashcard_set_id = f"flashcard_{ts}"
        if notebook_id:
            paths = get_notebook_paths(notebook_id, notebook_title or "", user_id)
            output_dir = paths.feature_output_dir("flashcard", ts)
        else:
            output_dir = _outputs_dir(email, notebook_id, flashcard_set_id)
        output_dir.mkdir(parents=True, exist_ok=True)

        flashcard_data = {
            "id": flashcard_set_id,
            "notebook_id": notebook_id,
            "flashcards": [fc.dict() for fc in flashcards],
            "created_at": time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime()),
            "source_files": file_paths,
            "total_count": len(flashcards),
        }
        (output_dir / "flashcards.json").write_text(
            json.dumps(flashcard_data, ensure_ascii=False, indent=2), encoding="utf-8"
        )
        log.info("[generate-flashcards] 成功生成 %d 张闪卡", len(flashcards))

        return {
            "success": True,
            "flashcards": [fc.dict() for fc in flashcards],
            "flashcard_set_id": flashcard_set_id,
            "total_count": len(flashcards),
            "result_path": _to_outputs_url(str(output_dir)),
        }
    except HTTPException:
        raise
    except Exception as e:
        log.exception("[generate-flashcards] failed")
        raise HTTPException(status_code=500, detail=str(e))


# ===================== Quiz 测验 =====================

@router.post("/generate-quiz")
async def generate_quiz(
    file_paths: List[str] = Body(..., embed=True),
    email: str = Body(..., embed=True),
    user_id: str = Body(..., embed=True),
    notebook_id: Optional[str] = Body(None, embed=True),
    notebook_title: Optional[str] = Body(None, embed=True),
    api_url: str = Body(..., embed=True),
    api_key: str = Body(..., embed=True),
    model: str = Body("deepseek-v3.2", embed=True),
    language: str = Body("en", embed=True),
    question_count: int = Body(10, embed=True),
):
    """生成 Quiz 测验题目"""
    try:
        from fastapi_app.services.quiz_service import generate_quiz_with_llm

        local_paths = []
        for f in file_paths:
            ps = (f or "").strip()
            if ps.startswith("http://") or ps.startswith("https://"):
                local_md = _resolve_link_to_local_md(email, notebook_id, ps)
                if local_md and local_md.exists():
                    local_paths.append(str(local_md))
            else:
                local_path = _resolve_local_path(f)
                if local_path.exists():
                    local_paths.append(str(local_path))

        if not local_paths:
            raise HTTPException(status_code=400, detail="No valid files provided")

        text_content = _extract_text_from_files(local_paths, max_chars=50000)
        if not text_content.strip():
            raise HTTPException(status_code=400, detail="No text content extracted")

        log.info("[generate-quiz] text_len=%d, files=%d", len(text_content), len(local_paths))

        questions = await generate_quiz_with_llm(
            text_content=text_content,
            api_url=api_url,
            api_key=api_key,
            model=model,
            language=language,
            question_count=question_count,
        )
        if not questions:
            raise HTTPException(status_code=500, detail="Failed to generate quiz")

        ts = int(time.time())
        quiz_id = f"quiz_{ts}"
        if notebook_id:
            paths = get_notebook_paths(notebook_id, notebook_title or "", user_id)
            output_dir = paths.feature_output_dir("quiz", ts)
        else:
            output_dir = _outputs_dir(email, notebook_id, quiz_id)
        output_dir.mkdir(parents=True, exist_ok=True)

        quiz_data = {
            "id": quiz_id,
            "notebook_id": notebook_id,
            "questions": [q.dict() for q in questions],
            "created_at": time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime()),
            "source_files": file_paths,
            "total_count": len(questions),
        }
        (output_dir / "quiz.json").write_text(
            json.dumps(quiz_data, ensure_ascii=False, indent=2), encoding="utf-8"
        )
        log.info("[generate-quiz] 成功生成 %d 道题目", len(questions))

        return {
            "success": True,
            "questions": [q.dict() for q in questions],
            "quiz_id": quiz_id,
            "total_count": len(questions),
            "result_path": _to_outputs_url(str(output_dir)),
        }
    except HTTPException:
        raise
    except Exception as e:
        log.exception("[generate-quiz] failed")
        raise HTTPException(status_code=500, detail=str(e))


# ===================== Flashcard / Quiz 读取端点 =====================

@router.get("/list-flashcard-sets")
async def list_flashcard_sets(
    notebook_id: str,
    notebook_title: Optional[str] = None,
    user_id: Optional[str] = None,
):
    """列出某 notebook 下所有已保存的闪卡集合（按时间倒序）"""
    try:
        paths = get_notebook_paths(notebook_id, notebook_title or "", user_id)
        flashcard_root = paths.root / "flashcard"
        sets = []
        if flashcard_root.exists():
            for ts_dir in flashcard_root.iterdir():
                if not ts_dir.is_dir():
                    continue
                json_file = ts_dir / "flashcards.json"
                if not json_file.exists():
                    continue
                try:
                    data = json.loads(json_file.read_text(encoding="utf-8"))
                    sets.append({
                        "set_id": ts_dir.name,
                        "id": data.get("id", ""),
                        "created_at": data.get("created_at", ""),
                        "total_count": data.get("total_count", 0),
                        "source_files": data.get("source_files", []),
                    })
                except Exception:
                    continue
        sets.sort(key=lambda x: x["set_id"], reverse=True)
        return {"success": True, "sets": sets}
    except Exception as e:
        log.exception("[list-flashcard-sets] failed")
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/list-quiz-sets")
async def list_quiz_sets(
    notebook_id: str,
    notebook_title: Optional[str] = None,
    user_id: Optional[str] = None,
):
    """列出某 notebook 下所有已保存的测验集合（按时间倒序）"""
    try:
        paths = get_notebook_paths(notebook_id, notebook_title or "", user_id)
        quiz_root = paths.root / "quiz"
        sets = []
        if quiz_root.exists():
            for ts_dir in quiz_root.iterdir():
                if not ts_dir.is_dir():
                    continue
                json_file = ts_dir / "quiz.json"
                if not json_file.exists():
                    continue
                try:
                    data = json.loads(json_file.read_text(encoding="utf-8"))
                    sets.append({
                        "set_id": ts_dir.name,
                        "id": data.get("id", ""),
                        "created_at": data.get("created_at", ""),
                        "total_count": data.get("total_count", 0),
                        "source_files": data.get("source_files", []),
                    })
                except Exception:
                    continue
        sets.sort(key=lambda x: x["set_id"], reverse=True)
        return {"success": True, "sets": sets}
    except Exception as e:
        log.exception("[list-quiz-sets] failed")
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/get-flashcard-set")
async def get_flashcard_set(
    notebook_id: str,
    set_id: str,
    notebook_title: Optional[str] = None,
    user_id: Optional[str] = None,
):
    """读取指定闪卡集合的完整数据"""
    try:
        paths = get_notebook_paths(notebook_id, notebook_title or "", user_id)
        json_file = paths.root / "flashcard" / set_id / "flashcards.json"
        if not json_file.exists():
            raise HTTPException(status_code=404, detail="Flashcard set not found")
        data = json.loads(json_file.read_text(encoding="utf-8"))
        return {"success": True, **data}
    except HTTPException:
        raise
    except Exception as e:
        log.exception("[get-flashcard-set] failed")
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/get-quiz-set")
async def get_quiz_set(
    notebook_id: str,
    set_id: str,
    notebook_title: Optional[str] = None,
    user_id: Optional[str] = None,
):
    """读取指定测验集合的完整数据"""
    try:
        paths = get_notebook_paths(notebook_id, notebook_title or "", user_id)
        json_file = paths.root / "quiz" / set_id / "quiz.json"
        if not json_file.exists():
            raise HTTPException(status_code=404, detail="Quiz set not found")
        data = json.loads(json_file.read_text(encoding="utf-8"))
        return {"success": True, **data}
    except HTTPException:
        raise
    except Exception as e:
        log.exception("[get-quiz-set] failed")
        raise HTTPException(status_code=500, detail=str(e))


# ============================================================================
# DeepResearch Integration
# ============================================================================

@router.post("/deep-research")
async def run_deep_research(
    query: str = Body(..., embed=True),
    notebook_id: str = Body(..., embed=True),
    notebook_title: Optional[str] = Body(None, embed=True),
    user_id: Optional[str] = Body(None, embed=True),
    email: Optional[str] = Body(None, embed=True),
    max_iterations: int = Body(50, embed=True),
):
    """
    运行 DeepResearch 深度研究并将结果保存为 source

    Args:
        query: 研究问题
        notebook_id: Notebook ID
        notebook_title: Notebook 标题
        user_id: 用户 ID
        email: 用户邮箱
        max_iterations: 最大迭代次数

    Returns:
        {
            "success": bool,
            "query": str,
            "answer": str,
            "source_info": {...},  # 保存的 source 信息
            "error": str (optional)
        }
    """
    try:
        from fastapi_app.services.deep_research_integration import DeepResearchIntegration

        log.info(f"[deep-research] 开始深度研究: {query}")

        # 1. 运行完整的 DeepResearch
        integration = DeepResearchIntegration()
        result = await integration.run_research(
            query=query,
            max_iterations=max_iterations
        )

        if not result["success"]:
            return result

        # 2. 将结果保存为 source
        paths = get_notebook_paths(notebook_id, notebook_title or "", user_id)
        mgr = SourceManager(paths)

        # 格式化为 Markdown
        markdown_content = integration.format_result_as_markdown(result)

        # 保存为文本 source
        source_info = await mgr.import_text(
            text=markdown_content,
            title=f"DeepResearch: {query[:50]}"
        )

        log.info(f"[deep-research] 已保存结果: {source_info.original_path}")

        # 3. 自动 embedding
        try:
            vector_base = str(paths.vector_store_dir)
            file_list = [{"path": str(source_info.original_path)}]
            await process_knowledge_base_files(
                file_list=file_list,
                vector_base=vector_base,
                email=email or "default",
                user_id=user_id or "default",
                notebook_id=notebook_id,
            )
            log.info(f"[deep-research] 已完成 embedding")
        except Exception as e:
            log.warning(f"[deep-research] Embedding 失败: {e}")

        return {
            "success": True,
            "query": query,
            "answer": result["answer"],
            "source_info": {
                "file_type": source_info.file_type,
                "original_path": str(source_info.original_path),
                "markdown_path": str(source_info.markdown_path) if source_info.markdown_path else None,
            },
            "sources_count": len(result.get("sources", [])),
        }

    except Exception as e:
        log.exception("[deep-research] 执行失败")
        raise HTTPException(status_code=500, detail=str(e))


# ============================================================================
# Search & Add Integration
# ============================================================================

@router.post("/search-and-add")
async def search_and_add(
    query: str = Body(..., embed=True),
    notebook_id: str = Body(..., embed=True),
    notebook_title: Optional[str] = Body(None, embed=True),
    user_id: Optional[str] = Body(None, embed=True),
    email: Optional[str] = Body(None, embed=True),
    top_k: int = Body(10, embed=True),
    search_provider: str = Body("serper", embed=True),
    search_api_key: Optional[str] = Body(None, embed=True),
):
    """
    搜索并爬取 Top K 结果，保存为 source

    Args:
        query: 搜索查询
        notebook_id: Notebook ID
        notebook_title: Notebook 标题
        user_id: 用户 ID
        email: 用户邮箱
        top_k: 返回前 K 个结果
        search_provider: 搜索引擎提供商
        search_api_key: 搜索 API 密钥

    Returns:
        {
            "success": bool,
            "query": str,
            "sources_count": int,
            "crawled_count": int,
            "source_info": {...}
        }
    """
    try:
        from fastapi_app.services.search_and_add_service import SearchAndAddService

        log.info(f"[search-and-add] 开始搜索: {query}, top_k={top_k}")

        # 1. 搜索并爬取
        service = SearchAndAddService()
        result = await service.search_and_crawl(
            query=query,
            top_k=top_k,
            search_provider=search_provider,
            search_api_key=search_api_key,
        )

        if not result["success"]:
            return result

        sources = result["sources"]
        if not sources:
            return {
                "success": False,
                "query": query,
                "error": "未找到搜索结果"
            }

        # 2. 将所有结果合并为一个 Markdown 文档
        paths = get_notebook_paths(notebook_id, notebook_title or "", user_id)
        mgr = SourceManager(paths)

        markdown_content = service.format_sources_as_markdown(sources)

        # 保存为文本 source
        source_info = await mgr.import_text(
            text=markdown_content,
            title=f"Search: {query[:50]}"
        )

        log.info(f"[search-and-add] 已保存 {len(sources)} 个结果: {source_info.original_path}")

        # 3. 自动 embedding
        try:
            vector_base = str(paths.vector_store_dir)
            file_list = [{"path": str(source_info.original_path)}]
            await process_knowledge_base_files(
                file_list=file_list,
                vector_base=vector_base,
                email=email or "default",
                user_id=user_id or "default",
                notebook_id=notebook_id,
            )
            log.info(f"[search-and-add] 已完成 embedding")
        except Exception as e:
            log.warning(f"[search-and-add] Embedding 失败: {e}")

        crawled_count = sum(1 for s in sources if s["crawl_success"])

        return {
            "success": True,
            "query": query,
            "sources_count": len(sources),
            "crawled_count": crawled_count,
            "source_info": {
                "file_type": source_info.file_type,
                "original_path": str(source_info.original_path),
                "markdown_path": str(source_info.markdown_path) if source_info.markdown_path else None,
            }
        }

    except Exception as e:
        log.exception("[search-and-add] 执行失败")
        raise HTTPException(status_code=500, detail=str(e))
