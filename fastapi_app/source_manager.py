"""
Source manager — unified import / read for notebook sources.

Handles:
- Copying originals into the new directory layout
- Running MinerU for PDFs
- Generating unified markdown for every source type
- Reading back markdown / MinerU data for feature cards
- Fallback to legacy kb_data / kb_mineru paths
- Structured chunk extraction with chunk_id / page_index / order / bbox (for GraphRAG)
"""
from __future__ import annotations

import asyncio
import hashlib
import json
import re
import shutil
import time
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from workflow_engine.logger import get_logger
from workflow_engine.utils import get_project_root

from fastapi_app.notebook_paths import NotebookPaths

log = get_logger(__name__)


@dataclass
class SourceInfo:
    stem: str
    original_path: Path
    markdown_path: Optional[Path] = None
    mineru_path: Optional[Path] = None
    file_type: str = ""  # pdf, md, docx, pptx, url, text


class SourceManager:
    """Manage sources for one notebook."""

    def __init__(self, paths: NotebookPaths):
        self.paths = paths

    # ------------------------------------------------------------------
    # Import
    # ------------------------------------------------------------------

    async def import_file(self, file_path: Path, filename: str) -> SourceInfo:
        """
        Import a local file into the notebook source tree.
        1. Copy to sources/{stem}/original/
        2. PDF → run MinerU → sources/{stem}/mineru/
        3. Generate unified markdown → sources/{stem}/markdown/{stem}.md
        """
        stem = Path(filename).stem
        ext = Path(filename).suffix.lower()

        # 1) Copy original
        orig_dir = self.paths.source_original_dir(filename)
        orig_dir.mkdir(parents=True, exist_ok=True)
        dest = orig_dir / filename
        if file_path.resolve() != dest.resolve():
            shutil.copy2(str(file_path), str(dest))

        info = SourceInfo(
            stem=stem,
            original_path=dest,
            file_type=ext.lstrip(".") or "unknown",
        )

        # 2) PDF → MinerU
        if ext == ".pdf":
            mineru_dir = self.paths.source_mineru_dir(filename)
            mineru_dir.mkdir(parents=True, exist_ok=True)
            try:
                await self._run_mineru(dest, mineru_dir)
                info.mineru_path = mineru_dir
            except Exception as e:
                log.warning("[SourceManager] MinerU failed for %s: %s", filename, e)

        # 3) Unified markdown
        md_dir = self.paths.source_markdown_dir(filename)
        md_dir.mkdir(parents=True, exist_ok=True)
        md_path = md_dir / f"{stem}.md"
        md_text = self._generate_markdown(dest, ext, info.mineru_path)
        if md_text:
            md_path.write_text(md_text, encoding="utf-8")
            info.markdown_path = md_path

        return info

    async def import_text(self, content: str, title: str) -> SourceInfo:
        """Import plain text as a .md source."""
        safe = re.sub(r'[^\w\u4e00-\u9fff\s\-.]', "", (title or "").strip())
        safe = (safe or "text")[:80].strip() or "text"
        filename = f"{safe}_{int(time.time())}.md"
        stem = Path(filename).stem

        orig_dir = self.paths.source_original_dir(filename)
        orig_dir.mkdir(parents=True, exist_ok=True)
        dest = orig_dir / filename
        dest.write_text((content or "").strip(), encoding="utf-8")

        md_dir = self.paths.source_markdown_dir(filename)
        md_dir.mkdir(parents=True, exist_ok=True)
        md_path = md_dir / f"{stem}.md"
        shutil.copy2(str(dest), str(md_path))

        return SourceInfo(
            stem=stem,
            original_path=dest,
            markdown_path=md_path,
            file_type="text",
        )

    async def import_url(self, url: str, fetched_text: str, title: str = "") -> SourceInfo:
        """Import a URL source (text already fetched by caller)."""
        if not title:
            from urllib.parse import urlparse
            parsed = urlparse(url)
            title = (parsed.netloc or "web") + "_" + (parsed.path.strip("/") or "page")[:30]
        safe = re.sub(r'[^\w\u4e00-\u9fff\s\-.]', "", title)
        safe = (safe or "url")[:80].strip() or "url"
        filename = f"{safe}_{int(time.time())}.md"
        stem = Path(filename).stem

        orig_dir = self.paths.source_original_dir(filename)
        orig_dir.mkdir(parents=True, exist_ok=True)
        dest = orig_dir / filename
        dest.write_text(fetched_text.strip(), encoding="utf-8")

        md_dir = self.paths.source_markdown_dir(filename)
        md_dir.mkdir(parents=True, exist_ok=True)
        md_path = md_dir / f"{stem}.md"
        shutil.copy2(str(dest), str(md_path))

        return SourceInfo(
            stem=stem,
            original_path=dest,
            markdown_path=md_path,
            file_type="url",
        )

    # ------------------------------------------------------------------
    # Read helpers
    # ------------------------------------------------------------------

    def get_markdown(self, source_stem: str) -> str:
        """Read unified markdown for a source. Falls back to legacy paths."""
        # New path
        md = self._find_in_sources(source_stem, "markdown", "*.md")
        if md:
            return md

        # Fallback: try reading original directly (for .md files)
        orig = self._find_in_sources(source_stem, "original", "*.md")
        if orig:
            return orig

        return ""

    def get_mineru_md(self, source_stem: str) -> str:
        """Read MinerU markdown for a PDF source."""
        md = self._find_in_sources(source_stem, "mineru", "*.md")
        if md:
            return md

        # Fallback: search inside mineru/{stem}/auto/*.md
        mineru_dir = self.paths.sources_dir / source_stem / "mineru"
        if mineru_dir.exists():
            for pattern in ["*/auto/*.md", "*/hybrid_auto/*.md", "*/*.md"]:
                for f in mineru_dir.glob(pattern):
                    try:
                        return f.read_text(encoding="utf-8")
                    except Exception:
                        continue
        return ""

    def get_mineru_root(self, source_stem: str) -> Optional[Path]:
        """Return the MinerU auto/ directory path (for images etc.)."""
        mineru_dir = self.paths.sources_dir / source_stem / "mineru"
        if not mineru_dir.exists():
            return None
        # Look for auto/ or hybrid_auto/
        for sub in ("auto", "hybrid_auto"):
            candidate = mineru_dir / sub
            if candidate.is_dir():
                return candidate
        # Fallback: {stem}/auto inside mineru dir
        for child in sorted(mineru_dir.iterdir()):
            if child.is_dir():
                for sub in ("auto", "hybrid_auto"):
                    c2 = child / sub
                    if c2.is_dir():
                        return c2
                # Any dir with .md files
                if list(child.glob("*.md")):
                    return child
        return None

    def get_sam3_cache_dir(self, source_stem: str) -> Optional[Path]:
        """Return the SAM3 cache directory if it exists."""
        sam3_dir = self.paths.sources_dir / source_stem / "sam3"
        return sam3_dir if sam3_dir.exists() else None

    def get_sam3_results(self, source_stem: str) -> Optional[dict]:
        """Read cached SAM3 results JSON for a source."""
        sam3_dir = self.get_sam3_cache_dir(source_stem)
        if not sam3_dir:
            return None
        json_path = sam3_dir / "sam3_results.json"
        if not json_path.exists():
            return None
        try:
            import json
            return json.loads(json_path.read_text(encoding="utf-8"))
        except Exception:
            return None

    def get_sam3_elements(self, source_stem: str) -> Optional[list]:
        """Read cached SAM3 drawio elements for a source."""
        sam3_dir = self.get_sam3_cache_dir(source_stem)
        if not sam3_dir:
            return None
        json_path = sam3_dir / "drawio_elements.json"
        if not json_path.exists():
            return None
        try:
            import json
            return json.loads(json_path.read_text(encoding="utf-8"))
        except Exception:
            return None

    def ensure_sam3_dir(self, source_stem: str) -> Path:
        """Create and return the SAM3 cache directory."""
        sam3_dir = self.paths.sources_dir / source_stem / "sam3"
        sam3_dir.mkdir(parents=True, exist_ok=True)
        return sam3_dir

    def get_chunks_with_meta(
        self,
        source_stem: str,
        chunk_size: int = 512,    # 默认值与 settings.GRAPHRAG_CHUNK_SIZE 一致
        chunk_overlap: int = 64,  # 默认值与 settings.GRAPHRAG_CHUNK_OVERLAP 一致
    ) -> List[Dict[str, Any]]:
        """Return structured chunks for a single source, used by GraphRAG indexing.

        Each dict has keys: chunk_id, text, page_index, order, bbox, source_stem.
        chunk_id = SHA1("{stem}:{order}")[:16], embedded as [chunk:ID] in input/*.txt.
        Priority: MinerU content_list.json (exact page+bbox) → MinerU MD (estimated page)
        → unified MD (page_index=-1).
        """
        chunks: List[Dict[str, Any]] = []

        # 1) MinerU content_list.json — exact page + bbox per block
        mineru_root = self.get_mineru_root(source_stem)
        if mineru_root:
            content_list_path = None
            # rglob to handle varying MinerU output directory layouts
            for candidate in mineru_root.parent.rglob("*_content_list.json"):
                content_list_path = candidate
                break
            if content_list_path and content_list_path.exists():
                try:
                    raw_blocks = json.loads(
                        content_list_path.read_text(encoding="utf-8")
                    )
                    order = 0
                    for block in raw_blocks:
                        # MinerU uses "text" or "content" depending on version
                        text = (block.get("text") or block.get("content") or "").strip()
                        if not text:
                            continue  # skip image / formula blocks
                        # MinerU uses "page_idx" or "page_index" depending on version
                        page_idx = int(block.get("page_idx", block.get("page_index", -1)))
                        bbox = block.get("bbox")  # [x1,y1,x2,y2] normalized, may be None
                        # chunk_id = SHA1("{stem}:{order}")[:16], embedded as [chunk:ID] in input/*.txt
                        chunk_id = hashlib.sha1(
                            f"{source_stem}:{order}".encode()
                        ).hexdigest()[:16]
                        chunks.append(
                            {
                                "chunk_id": chunk_id,
                                "text": text,
                                "page_index": page_idx,
                                "order": order,
                                "bbox": bbox,
                                "source_stem": source_stem,
                            }
                        )
                        order += 1
                    if chunks:
                        return chunks
                except Exception as e:
                    log.debug(
                        "[SourceManager] content_list.json parse failed for %s: %s",
                        source_stem,
                        e,
                    )

        # 2) MinerU markdown fallback — sliding window, estimated page_index
        mineru_md = self.get_mineru_md(source_stem)
        if mineru_md.strip():
            chunks = self._split_text_to_chunks(
                mineru_md, source_stem, chunk_size, chunk_overlap, estimate_pages=True
            )
            if chunks:
                return chunks

        # 3) Unified markdown fallback — no page info (Word/PPT/TXT)
        md = self.get_markdown(source_stem)
        if md.strip():
            return self._split_text_to_chunks(
                md, source_stem, chunk_size, chunk_overlap, estimate_pages=False
            )

        return []

    def get_all_markdowns(self) -> List[Tuple[str, str]]:
        """Return [(stem, markdown_text), ...] for all sources."""
        results: List[Tuple[str, str]] = []
        sources_dir = self.paths.sources_dir
        if not sources_dir.exists():
            return results
        for src_dir in sorted(sources_dir.iterdir()):
            if not src_dir.is_dir():
                continue
            stem = src_dir.name
            md = self.get_markdown(stem)
            if md.strip():
                results.append((stem, md))
        return results

    def get_original_path(self, source_stem: str) -> Optional[Path]:
        """Return the original file path for a source."""
        orig_dir = self.paths.sources_dir / source_stem / "original"
        if not orig_dir.exists():
            return None
        for f in orig_dir.iterdir():
            if f.is_file():
                return f
        return None

    def list_sources(self) -> List[SourceInfo]:
        """List all sources in this notebook."""
        results: List[SourceInfo] = []
        sources_dir = self.paths.sources_dir
        if not sources_dir.exists():
            return results
        for src_dir in sorted(sources_dir.iterdir()):
            if not src_dir.is_dir():
                continue
            stem = src_dir.name
            orig = self.get_original_path(stem)
            if not orig:
                continue
            ext = orig.suffix.lower().lstrip(".")
            md_dir = src_dir / "markdown"
            md_path = None
            if md_dir.exists():
                mds = list(md_dir.glob("*.md"))
                if mds:
                    md_path = mds[0]
            mineru_path = src_dir / "mineru" if (src_dir / "mineru").exists() else None
            results.append(SourceInfo(
                stem=stem,
                original_path=orig,
                markdown_path=md_path,
                mineru_path=mineru_path,
                file_type=ext,
            ))
        return results

    # ------------------------------------------------------------------
    # Internal
    # ------------------------------------------------------------------

    async def _run_mineru(self, pdf_path: Path, output_dir: Path) -> None:
        """Run MinerU on a PDF file."""
        from workflow_engine.toolkits.multimodaltool.mineru_tool import run_mineru_pdf_extract
        await asyncio.to_thread(
            run_mineru_pdf_extract,
            str(pdf_path),
            str(output_dir),
            "modelscope",
            None,
            "pipeline",
        )

    def _generate_markdown(
        self, file_path: Path, ext: str, mineru_dir: Optional[Path]
    ) -> str:
        """Generate unified markdown from a source file."""
        # PDF: copy MinerU's .md output
        if ext == ".pdf" and mineru_dir:
            for pattern in ["*.md", "*/auto/*.md", "*/hybrid_auto/*.md"]:
                for f in mineru_dir.glob(pattern):
                    try:
                        return f.read_text(encoding="utf-8")
                    except Exception:
                        continue
            # Fallback: PyMuPDF extraction
            return self._extract_text_pymupdf(file_path)

        if ext == ".pdf":
            return self._extract_text_pymupdf(file_path)

        # MD / TXT: direct copy
        if ext in (".md", ".markdown", ".txt"):
            try:
                return file_path.read_text(encoding="utf-8", errors="replace")
            except Exception:
                return ""

        # DOCX
        if ext in (".docx", ".doc"):
            return self._extract_text_docx(file_path)

        # PPTX
        if ext in (".pptx", ".ppt"):
            return self._extract_text_pptx(file_path)

        # Fallback: try reading as text
        try:
            return file_path.read_text(encoding="utf-8", errors="replace")
        except Exception:
            return ""

    def _extract_text_pymupdf(self, path: Path) -> str:
        try:
            import fitz
            doc = fitz.open(str(path))
            text = "\n\n".join(page.get_text() for page in doc)
            doc.close()
            return text
        except Exception as e:
            log.warning("[SourceManager] PyMuPDF extraction failed: %s", e)
            return ""

    def _extract_text_docx(self, path: Path) -> str:
        try:
            from docx import Document
            doc = Document(str(path))
            return "\n".join(p.text for p in doc.paragraphs)
        except Exception as e:
            log.warning("[SourceManager] docx extraction failed: %s", e)
            return ""

    def _extract_text_pptx(self, path: Path) -> str:
        try:
            from pptx import Presentation
            prs = Presentation(str(path))
            parts = []
            for i, slide in enumerate(prs.slides):
                parts.append(f"--- Slide {i+1} ---")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        parts.append(shape.text)
            return "\n".join(parts)
        except Exception as e:
            log.warning("[SourceManager] pptx extraction failed: %s", e)
            return ""

    def _find_in_sources(self, source_stem: str, subdir: str, pattern: str) -> str:
        """Find and read the first matching file in sources/{stem}/{subdir}/."""
        d = self.paths.sources_dir / source_stem / subdir
        if not d.exists():
            return ""
        for f in d.glob(pattern):
            try:
                return f.read_text(encoding="utf-8", errors="replace")
            except Exception:
                continue
        return ""

    @staticmethod
    def _split_text_to_chunks(
        text: str,
        source_stem: str,
        chunk_size: int,
        chunk_overlap: int,
        estimate_pages: bool,
    ) -> List[Dict[str, Any]]:
        """Sliding-window character chunking fallback when content_list is unavailable.

        estimate_pages=True roughly estimates page_index at ~2000 chars/page.
        """
        chunks: List[Dict[str, Any]] = []
        text = text.strip()
        if not text:
            return chunks

        total_chars = len(text)
        step = max(1, chunk_size - chunk_overlap)
        order = 0
        pos = 0
        chars_per_page = 2000  # rough estimate: ~2000 chars per page

        while pos < total_chars:
            end = min(pos + chunk_size, total_chars)
            snippet = text[pos:end].strip()
            if snippet:
                page_idx = int(pos / chars_per_page) if estimate_pages else -1
                chunk_id = hashlib.sha1(
                    f"{source_stem}:{order}".encode()
                ).hexdigest()[:16]
                chunks.append(
                    {
                        "chunk_id": chunk_id,
                        "text": snippet,
                        "page_index": page_idx,
                        "order": order,
                        "bbox": None,
                        "source_stem": source_stem,
                    }
                )
                order += 1
            pos += step

        return chunks
